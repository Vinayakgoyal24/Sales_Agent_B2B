import os
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from io import BytesIO
from dotenv import load_dotenv

# Load environment variables
load_dotenv()
import os
from dotenv import load_dotenv
load_dotenv()
print("🔐 Key:", os.getenv("SERPAPI_KEY"))


def search_image_url_serpapi(product_name: str) -> str:
    serp_api_key = os.getenv("SERPAPI_KEY")
    if not serp_api_key:
        print("⚠️ SERPAPI_KEY missing in environment.")
        return ""

    params = {
        "engine": "google_images",
        "q": product_name,
        "api_key": serp_api_key,
        "num": 1,
        "safe": "active",
        "hl": "en",
    }

    try:
        response = requests.get("https://serpapi.com/search.json", params=params)
        response.raise_for_status()
        data = response.json()
        if data.get("images_results"):
            return data["images_results"][0]["original"]
    except Exception as e:
        print(f"[Image Search Error] {e}")
    return ""

def download_image(image_url: str) -> BytesIO or None:
    try:
        response = requests.get(image_url, timeout=5)
        response.raise_for_status()
        return BytesIO(response.content)
    except Exception as e:
        print(f"[Image Download Error] {e}")
        return None

def create_test_slide(product_name: str):
    print(f"🔍 Searching image for: {product_name}")
    image_url = search_image_url_serpapi(product_name)
    print("📷 Image URL:", image_url)

    if not image_url:
        print("❌ No image URL found.")
        return

    image_data = download_image(image_url)
    if not image_data:
        print("❌ Failed to download image.")
        return

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide_width, slide_height = prs.slide_width, prs.slide_height

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    p = title_box.text_frame.add_paragraph()
    p.text = f"Product: {product_name}"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    # Add Image
    try:
        slide.shapes.add_picture(image_data, Inches(1), Inches(1.5), width=Inches(4.5))
        print("✅ Image added to slide.")
    except Exception as e:
        print(f"❌ Error inserting image: {e}")

    prs.save("test_slide.pptx")
    print("📄 Saved to test_slide.pptx")

# 🔧 Run test
create_test_slide("ThinkPad X1 Carbon")
