from fastapi import FastAPI, HTTPException, File, UploadFile
from fastapi.responses import StreamingResponse, JSONResponse
from pydantic import BaseModel
from typing import List, Optional
import os
import pandas as pd
from dotenv import load_dotenv
from langchain_community.document_loaders import TextLoader
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_chroma import Chroma
from langchain_openai import AzureOpenAIEmbeddings, AzureChatOpenAI
from langchain.prompts import ChatPromptTemplate
from langgraph.graph import START, StateGraph
from typing_extensions import TypedDict
import tiktoken
from io import BytesIO
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER
from reportlab.pdfbase.pdfmetrics import stringWidth
from datetime import datetime, timedelta
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import logging


from fastapi import FastAPI
from fastapi.middleware.cors import CORSMiddleware

app = FastAPI()

# Add CORS middleware BEFORE defining your routes
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Allow all origins for development
    allow_credentials=True,
    allow_methods=["GET", "POST", "OPTIONS"],  # Explicitly include OPTIONS
    allow_headers=["*"],
)

@app.post("/query")
async def query_endpoint(request_data: dict):
    # Your existing code here
    pass



# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Load environment variables
load_dotenv()

# Initialize FastAPI app
app = FastAPI(
    title="Hardware Sales Assistant API",
    description="AI-powered hardware quotation and recommendation system",
    version="1.0.0"
)

# --- Pydantic Models ---
class QueryRequest(BaseModel):
    question: str
    
class QueryResponse(BaseModel):
    answer: str
    input_tokens: int
    output_tokens: int
    total_tokens: int
    context_documents: int

class HealthResponse(BaseModel):
    status: str
    message: str

class DocumentUploadResponse(BaseModel):
    message: str
    documents_processed: int

# --- Initialize Azure Clients ---
try:
    llm = AzureChatOpenAI(
        azure_endpoint=os.getenv("AZURE_OPENAI_ENDPOINT"),
        azure_deployment=os.getenv("AZURE_OPENAI_DEPLOYMENT_NAME"),
        openai_api_version=os.getenv("AZURE_OPENAI_API_VERSION"),
    )

    embeddings = AzureOpenAIEmbeddings(
        azure_endpoint=os.getenv("AZURE_OPENAI_EMBEDDING_ENDPOINT"),
        azure_deployment=os.getenv("AZURE_OPENAI_EMBEDDING_DEPLOYMENT_NAME"),
        openai_api_version=os.getenv("AZURE_OPENAI_EMBEDDING_API_VERSION"),
    )

    vector_store = Chroma(
        collection_name="example_collection",
        embedding_function=embeddings,
        persist_directory="./chroma_langchain_db",
    )
    
    logger.info("Azure OpenAI clients initialized successfully")
except Exception as e:
    logger.error(f"Failed to initialize Azure OpenAI clients: {e}")
    raise

# --- Prompt Template ---
prompt = ChatPromptTemplate.from_messages([
    ("system",
     "You are a professional hardware sales assistant at Otsuka Corporation. Based on the user's request and the context, provide 2–3 detailed hardware configuration quotations. "
     "Each quotation should include:\n"
     "- Product Name\n- Specs\n- Price\n- Quantity\n- Total Price\n\n"
     "Use this structure:\n"
     "## Quotation 1\nProduct Name: ...\nSpecs: ...\nPrice: ...\nQuantity: ...\n...\nTotal Price: ...\n"
     "## Quotation 2 ...\n\n"
     "Then provide a clear comparison of the quotations and recommend the best one based on:\n"
     "- Price\n- Suitability for the user's need\n- Performance vs cost.\n"
     "Use a section titled:\n"
     "## Recommendation\n"
     "Mention why the chosen quote is the best and highlight key differences with others.\n\n"
     "Keep tone professional and brief. Do not fabricate information if context is insufficient."),
    ("human", "Question: {question}\n\nContext:\n{context}")
])

# --- Helper Functions ---
def load_csv_as_documents(folder_path="data"):
    """Load CSV files as documents"""
    documents = []
    if not os.path.exists(folder_path):
        logger.warning(f"Data folder {folder_path} does not exist")
        return documents
        
    for file_name in os.listdir(folder_path):
        if file_name.endswith(".csv"):
            try:
                df = pd.read_csv(os.path.join(folder_path, file_name)).head(20)
                for _, row in df.iterrows():
                    content = "\n".join([f"{k}: {v}" for k, v in row.items()])
                    documents.append(Document(page_content=content, metadata={"source": file_name}))
                logger.info(f"Loaded {len(df)} rows from {file_name}")
            except Exception as e:
                logger.error(f"Error loading {file_name}: {e}")
    return documents

def initialize_vector_store():
    """Initialize vector store with CSV data"""
    try:
        if vector_store._collection.count() == 0:
            docs = load_csv_as_documents("data")
            if docs:
                splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
                splits = splitter.split_documents(docs)
                vector_store.add_documents(splits)
                vector_store.persist()
                logger.info(f"Indexed {len(splits)} document chunks")
            else:
                logger.warning("No documents found to index")
        else:
            logger.info(f"Vector store already contains {vector_store._collection.count()} documents")
    except Exception as e:
        logger.error(f"Error initializing vector store: {e}")
        raise

def generate_pdf(quotation_text: str) -> BytesIO:
    """Generate PDF from quotation text"""
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer,
                            pagesize=A4,
                            rightMargin=30, leftMargin=30,
                            topMargin=30, bottomMargin=30)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(name="TitleCenter", parent=styles["Title"], alignment=TA_CENTER)
    highlight_style = ParagraphStyle(name="Highlight", parent=styles["Normal"], textColor=colors.green, fontSize=12)

    elements = []

    # Header: Company Info
    logo_path = "otsuka_im.png"
    if os.path.exists(logo_path):
        logo = Image(logo_path, width=100, height=50)
        logo.hAlign = "LEFT"
        elements.append(logo)

    elements.append(Paragraph("Otsuka Corporation", title_style))
    elements.append(Paragraph("Head Office, 2-18-4 Iidabashi, Chiyoda-ku, Tokyo 102-8573", styles["Normal"]))
    elements.append(Paragraph("Website: <a href='https://www.otsuka-shokai.co.jp/'>www.otsuka-shokai.co.jp</a>",
                              styles["Normal"]))
    elements.append(Spacer(1, 12))

    # Quotation Info
    elements.append(Paragraph("<b>🧾 Quotation</b>", styles["Heading2"]))
    today = datetime.now()
    validity = today + timedelta(days=7)
    elements.append(Paragraph(f"Date of Issue: {today.strftime('%Y-%m-%d')}", styles["Normal"]))
    elements.append(Paragraph(f"Validity: {validity.strftime('%Y-%m-%d')} (7 days)", styles["Normal"]))
    elements.append(Spacer(1, 12))

    # Client Info
    elements.append(Paragraph("<b>📋 Client Information</b>", styles["Heading3"]))
    client_info = [
        "Client Name: Acme Solutions Pvt. Ltd.",
        "Client Address: 123, Innovation Tower, Marunouchi, Tokyo",
        "Contact Person: John Doe",
        "Email: john.doe@example.com",
        "Phone: +81 90-1234-5678"
    ]
    for line in client_info:
        elements.append(Paragraph(line, styles["Normal"]))
    elements.append(Spacer(1, 12))

    # Parse quotation text and build tables
    lines = quotation_text.strip().splitlines()
    table_data = []
    total_prices = []
    price_qty_list = []
    recommendation_lines = []
    current_quotation = ""
    inside_quote = False

    def build_table(title, data, bg_color):
        table_style = styles
        cell_style = styles["Normal"]
        font_name = "Helvetica"
        font_size = 10

        tbl = [data[0]]
        for row in data[1:]:
            tbl.append([Paragraph(str(cell), cell_style) for cell in row])

        transposed = list(zip(*data))
        col_widths = []
        for col in transposed:
            max_w = max(stringWidth(str(item), font_name, font_size) for item in col)
            col_widths.append(min(max_w + 20, 200))

        table = Table(tbl, hAlign="LEFT", colWidths=col_widths)
        table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), bg_color),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.whitesmoke),
            ("GRID", (0, 0), (-1, -1), 0.7, colors.black),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), font_size),
            ("ALIGN", (2, 1), (2, -1), "RIGHT"),
            ("ALIGN", (3, 1), (3, -1), "CENTER"),
            ("VALIGN", (0, 1), (-1, -1), "TOP"),
        ]))

        elements.append(Paragraph(f"<b>{title}</b>", styles["Heading4"]))
        elements.append(table)

    for line in lines:
        line = line.strip()

        if line.startswith("## Quotation"):
            if table_data:
                subtotal = sum(p * q for p, q in price_qty_list)
                subtotal_str = f"{subtotal:,.0f}"
                build_table(current_quotation, table_data, bg_color=colors.HexColor("#4472C4"))
                elements.append(Paragraph(f"<b>Total Price ({current_quotation}):</b> ¥{subtotal_str}", styles["Normal"]))
                total_prices.append((current_quotation, subtotal_str))
                elements.append(Spacer(1, 10))

            current_quotation = line.replace("##", "").strip()
            table_data = [["Product Name", "Specs", "Price ($)", "Qty"]]
            price_qty_list = []
            inside_quote = True

        elif line.startswith("Product Name:"):
            pname = line.split(":", 1)[1].strip()

        elif line.startswith("Specs:"):
            specs = line.split(":", 1)[1].strip()

        elif line.startswith("Price:"):
            raw = line.split(":", 1)[1].strip()
            clean = re.sub(r"[^\d\.]", "", raw)
            price = float(clean)

        elif line.startswith("Quantity:"):
            qty = int(line.split(":", 1)[1].strip())
            table_data.append([pname, specs, f"{price:,.0f}", str(qty)])
            price_qty_list.append((price, qty))

        elif line.startswith("## Recommendation"):
            if table_data and inside_quote:
                subtotal = sum(p * q for p, q in price_qty_list)
                subtotal_str = f"{subtotal:,.0f}"
                build_table(current_quotation, table_data, bg_color=colors.HexColor("#4472C4"))
                elements.append(Paragraph(f"<b>Total Price ({current_quotation}):</b> ¥{subtotal_str}", styles["Normal"]))
                total_prices.append((current_quotation, subtotal_str))
                elements.append(Spacer(1, 10))
            inside_quote = False
            elements.append(Paragraph("<b>🎯 Recommendation</b>", styles["Heading3"]))

        elif not inside_quote and line:
            recommendation_lines.append(line)

    if table_data and inside_quote:
        build_table(current_quotation, table_data, bg_color=colors.HexColor("#4472C4"))

    # Pricing Summary
    elements.append(Spacer(1, 12))
    elements.append(Paragraph("<b>📊 Pricing Summary</b>", styles["Heading3"]))
    for qname, tprice in total_prices:
        elements.append(Paragraph(f"• {qname}: ¥{tprice}", styles["Normal"]))

    # Recommendation Section
    if recommendation_lines:
        elements.append(Spacer(1, 12))
        elements.append(Paragraph("<b>✅ Best Recommendation</b>", styles["Heading3"]))
        for line in recommendation_lines:
            elements.append(Paragraph(line, highlight_style))

    doc.build(elements)
    buffer.seek(0)
    return buffer

def generate_slides(quotation_text: str) -> BytesIO:
    """Generate PowerPoint slides from quotation text"""
    prs = Presentation()
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    def add_logo(slide):
        if os.path.exists("otsuka_im.png"):
            LOGO_WIDTH_INCHES = 1.2
            LOGO_TOP_MARGIN = 0.2
            LOGO_RIGHT_MARGIN = 0.3
            left = slide_width - Inches(LOGO_RIGHT_MARGIN + LOGO_WIDTH_INCHES)
            top = Inches(LOGO_TOP_MARGIN)
            slide.shapes.add_picture("otsuka_im.png", left, top, width=Inches(LOGO_WIDTH_INCHES))

    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Hardware Configuration Quotations"
    slide.placeholders[1].text = "Generated by Otsuka Corporation's Sales Assistant"
    add_logo(slide)

    # Client Info Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "📋 Client Information"
    content = slide.placeholders[1]
    content.text = "\n".join([
        "Client Name: Acme Solutions Pvt. Ltd.",
        "Client Address: 123, Innovation Tower, Marunouchi, Tokyo",
        "Contact Person: John Doe",
        "Email: john.doe@example.com",
        "Phone: +81 90-1234-5678"
    ])
    add_logo(slide)

    lines = quotation_text.strip().splitlines()
    current_quotation = ""
    recommendation_lines = []
    table_data = []
    price_qty_list = []
    inside_quote = False

    def add_quotation_slide(prs, title, table_data, price_qty_list):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        shapes = slide.shapes
        shapes.title.text = title

        rows = len(table_data)
        cols = len(table_data[0])
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(0.8 + 0.4 * rows)

        table = shapes.add_table(rows, cols, left, top, width, height).table

        for col_index in range(cols):
            table.columns[col_index].width = Inches(2.2)

        for r in range(rows):
            for c in range(cols):
                cell = table.cell(r, c)
                cell.text = str(table_data[r][c])
                cell.text_frame.paragraphs[0].font.size = Pt(11)

        total = sum(p * q for p, q in price_qty_list)
        txBox = slide.shapes.add_textbox(Inches(0.5), top + height + Inches(0.3), Inches(5), Inches(1))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = f"Total Price: ¥{total:,.0f}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 112, 192)
        add_logo(slide)

    for line in lines:
        line = line.strip()

        if line.startswith("## Quotation"):
            if table_data:
                add_quotation_slide(prs, current_quotation, table_data, price_qty_list)
            current_quotation = line.replace("##", "").strip()
            table_data = [["Product Name", "Specs", "Price", "Qty"]]
            price_qty_list = []
            inside_quote = True

        elif line.startswith("Product Name:"):
            pname = line.split(":", 1)[1].strip()

        elif line.startswith("Specs:"):
            specs = line.split(":", 1)[1].strip()

        elif line.startswith("Price:"):
            raw = line.split(":", 1)[1].strip()
            clean = re.sub(r"[^\d\.]", "", raw)
            price = float(clean)

        elif line.startswith("Quantity:"):
            qty = int(line.split(":", 1)[1].strip())
            table_data.append([pname, specs, f"${price:,.0f}", str(qty)])
            price_qty_list.append((price, qty))

        elif line.startswith("## Recommendation"):
            if table_data:
                add_quotation_slide(prs, current_quotation, table_data, price_qty_list)
                table_data = []
            inside_quote = False

        elif not inside_quote and line:
            recommendation_lines.append(line)

    if table_data and inside_quote:
        add_quotation_slide(prs, current_quotation, table_data, price_qty_list)

    # Recommendation Slide
    if recommendation_lines:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "🎯 Recommendation"
        content = slide.placeholders[1]
        content.text = "\n".join(recommendation_lines)
        add_logo(slide)

    # Thank You Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "🙏 Thank You"
    content = slide.placeholders[1]
    content.text = "\n".join([
        "Otsuka Corporation",
        "Head Office: 2-18-4 Iidabashi, Chiyoda-ku, Tokyo 102-8573",
        "Website: https://www.otsuka-shokai.co.jp"
    ])
    add_logo(slide)

    pptx_buffer = BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    return pptx_buffer

# --- LangGraph Setup ---
class State(TypedDict):
    question: str
    context: List[Document]
    answer: str

def retrieve(state: State):
    return {"context": vector_store.similarity_search(state["question"])}

def generate(state: State):
    docs_content = "\n\n".join(doc.page_content for doc in state["context"])
    messages = prompt.invoke({"question": state["question"], "context": docs_content})
    response = llm.invoke(messages)
    return {"answer": response.content}

graph_builder = StateGraph(State).add_sequence([retrieve, generate])
graph_builder.add_edge(START, "retrieve")
graph = graph_builder.compile()

# Initialize vector store on startup
initialize_vector_store()

# --- API Endpoints ---

@app.get("/", response_model=HealthResponse)
async def root():
    """Health check endpoint"""
    return HealthResponse(status="healthy", message="Hardware Sales Assistant API is running")

@app.get("/health", response_model=HealthResponse)
async def health_check():
    """Detailed health check"""
    try:
        # Test vector store
        count = vector_store._collection.count()
        return HealthResponse(
            status="healthy", 
            message=f"API is running. Vector store contains {count} documents."
        )
    except Exception as e:
        return HealthResponse(status="unhealthy", message=f"Error: {str(e)}")

@app.post("/query", response_model=QueryResponse)
async def process_query(request: QueryRequest):
    """Process hardware recommendation query"""
    try:
        if not request.question.strip():
            raise HTTPException(status_code=400, detail="Question cannot be empty")
        
        logger.info(f"Processing query: {request.question[:100]}...")
        
        # Process query through LangGraph
        result = graph.invoke({"question": request.question})
        
        # Calculate token usage
        enc = tiktoken.encoding_for_model("gpt-4")
        input_tokens = len(enc.encode(request.question))
        output_tokens = len(enc.encode(result["answer"]))
        context_docs = len(result.get("context", []))
        
        logger.info(f"Query processed successfully. Tokens: {input_tokens + output_tokens}")
        
        return QueryResponse(
            answer=result["answer"],
            input_tokens=input_tokens,
            output_tokens=output_tokens,
            total_tokens=input_tokens + output_tokens,
            context_documents=context_docs
        )
        
    except Exception as e:
        logger.error(f"Error processing query: {e}")
        raise HTTPException(status_code=500, detail=f"Error processing query: {str(e)}")

@app.post("/generate-pdf")
async def generate_pdf_endpoint(request: QueryRequest):
    """Generate PDF quotation from query"""
    try:
        if not request.question.strip():
            raise HTTPException(status_code=400, detail="Question cannot be empty")
        
        logger.info(f"Generating PDF for query: {request.question[:100]}...")
        
        # Get quotation
        result = graph.invoke({"question": request.question})
        
        # Generate PDF
        pdf_buffer = generate_pdf(result["answer"])
        
        return StreamingResponse(
            BytesIO(pdf_buffer.read()),
            media_type="application/pdf",
            headers={"Content-Disposition": "attachment; filename=hardware_quotation.pdf"}
        )
        
    except Exception as e:
        logger.error(f"Error generating PDF: {e}")
        raise HTTPException(status_code=500, detail=f"Error generating PDF: {str(e)}")

@app.post("/generate-slides")
async def generate_slides_endpoint(request: QueryRequest):
    """Generate PowerPoint slides from query"""
    try:
        if not request.question.strip():
            raise HTTPException(status_code=400, detail="Question cannot be empty")
        
        logger.info(f"Generating slides for query: {request.question[:100]}...")
        
        # Get quotation
        result = graph.invoke({"question": request.question})
        
        # Generate slides
        slides_buffer = generate_slides(result["answer"])
        
        return StreamingResponse(
            BytesIO(slides_buffer.read()),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": "attachment; filename=hardware_quotation.pptx"}
        )
        
    except Exception as e:
        logger.error(f"Error generating slides: {e}")
        raise HTTPException(status_code=500, detail=f"Error generating slides: {str(e)}")

@app.post("/upload-csv", response_model=DocumentUploadResponse)
async def upload_csv(file: UploadFile = File(...)):
    """Upload and index a CSV file"""
    try:
        if not file.filename.endswith('.csv'):
            raise HTTPException(status_code=400, detail="Only CSV files are allowed")
        
        # Save uploaded file
        os.makedirs("data", exist_ok=True)
        file_path = os.path.join("data", file.filename)
        
        with open(file_path, "wb") as f:
            content = await file.read()
            f.write(content)
        
        # Process and index the new file
        df = pd.read_csv(file_path).head(20)
        documents = []
        
        for _, row in df.iterrows():
            content = "\n".join([f"{k}: {v}" for k, v in row.items()])
            documents.append(Document(page_content=content, metadata={"source": file.filename}))
        
        # Add to vector store
        splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        splits = splitter.split_documents(documents)
        vector_store.add_documents(splits)
        vector_store.persist()
        
        logger.info(f"Uploaded and indexed {file.filename} with {len(documents)} rows")
        
        return DocumentUploadResponse(
            message=f"Successfully uploaded and indexed {file.filename}",
            documents_processed=len(documents)
        )
        
    except Exception as e:
        logger.error(f"Error uploading CSV: {e}")
        raise HTTPException(status_code=500, detail=f"Error uploading CSV: {str(e)}")

@app.get("/vector-store/stats")
async def get_vector_store_stats():
    """Get vector store statistics"""
    try:
        count = vector_store._collection.count()
        return {
            "total_documents": count,
            "collection_name": vector_store._collection.name,
            "status": "active"
        }
    except Exception as e:
        logger.error(f"Error getting vector store stats: {e}")
        raise HTTPException(status_code=500, detail=f"Error getting stats: {str(e)}")

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)