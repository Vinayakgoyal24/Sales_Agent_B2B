import os
import re
import pandas as pd
import tiktoken
from dotenv import load_dotenv
from io import BytesIO
import requests
import time
import streamlit as st
from langchain_community.document_loaders import TextLoader
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_chroma import Chroma
from langchain_openai import AzureOpenAIEmbeddings, AzureChatOpenAI
from langchain.prompts import ChatPromptTemplate
from langgraph.graph import START, StateGraph
from typing_extensions import List, TypedDict
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER
from reportlab.pdfbase.pdfmetrics import stringWidth
from datetime import datetime, timedelta
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from PIL import Image as PILImage
import io
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from email.mime.base import MIMEBase
from email import encoders
import pandas as pd
import base64
import tiktoken
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from io import BytesIO
import os
import re


def search_image_url_serpapi(product_name: str) -> str:
    print(f"🔍 Searching image for: {product_name}")
    serp_api_key = os.getenv("SERPAPI_KEY")
    if not serp_api_key:
        print("⚠️ SERPAPI_KEY missing in environment. Please set it in your .env file.")
        return ""
    params = {
        "engine": "google_images",
        "q": product_name,
        "api_key": serp_api_key,
        "num": 1,
        "safe": "active",
        "hl": "en",
    }
    try:
        response = requests.get("https://serpapi.com/search.json", params=params)
        response.raise_for_status()
        data = response.json()
        if data.get("images_results"):
            image_url = data["images_results"][0]["original"]
            print(f"✅ Found image: {image_url}")
            return image_url
        else:
            print(f"❌ No image results found for: {product_name}")
    except Exception as e:
        print(f"[Image Search Error for '{product_name}']: {e}")
    return ""


def download_image(image_url: str) -> BytesIO or None:
    try:
        start = time.time()
        response = requests.get(image_url, timeout=5)
        response.raise_for_status()
        print(f"📥 Downloaded image in {time.time() - start:.2f} sec")
        return BytesIO(response.content)
    except Exception as e:
        print(f"[Image Download Error] {e}")
        return None


def contains_japanese(text: str) -> bool:
    return bool(re.search(r"[\u3040-\u30FF\u4E00-\u9FFF]", text))


def generate_slides(quotation_text: str, client_info: dict) -> BytesIO:



    starttimeppt = time.time()
    prs = Presentation()
    slide_width, slide_height = prs.slide_width, prs.slide_height
    bg_color = RGBColor(245, 245, 245)

    def set_slide_bg_color(slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    def add_logo(slide):
        if os.path.exists("otsuka_im.png"):
            left, top = slide_width - Inches(1.5), Inches(0.3)
            slide.shapes.add_picture("otsuka_im.png", left, top, width=Inches(1.2))

    def remove_title_placeholder(slide):
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                sp = shape
                sp.element.getparent().remove(sp.element)

    def add_footer(slide, index):
        margin_bottom = Inches(0.2)
        textbox_height = Inches(0.25)

        footer = slide.shapes.add_textbox(
            Inches(0.3), slide_height - margin_bottom - textbox_height, Inches(3), textbox_height
        )
        p = footer.text_frame.add_paragraph()
        p.text = "by Otsuka Shokai"
        p.font.size = Pt(10)
        p.font.italic = True
        p.font.color.rgb = RGBColor(120, 120, 120)

        page_num = slide.shapes.add_textbox(
            slide_width - Inches(1), slide_height - margin_bottom - textbox_height, Inches(0.7), textbox_height
        )
        p2 = page_num.text_frame.add_paragraph()
        p2.text = f"Page {index}"
        p2.font.size = Pt(10)
        p2.font.color.rgb = RGBColor(120, 120, 120)

    def underline_title(slide):
        if slide.shapes.title:
            title_shape = slide.shapes.title
            p = title_shape.text_frame.paragraphs[0]
            p.font.underline = True
            p.font.color.rgb = RGBColor(0, 51, 102)
            p.font.size = Pt(28)

    def add_quotation_slide(prs, title, table_data, price_qty_list, index):
        print(f"\n📄 Creating slide: {title}")
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        set_slide_bg_color(slide)
        slide.shapes.title.text = title
        underline_title(slide)

        rows, cols = len(table_data), len(table_data[0])
        left, top, width = Inches(0.5), Inches(1.5), Inches(9)
        height = Inches(0.8 + 0.4 * rows)
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        for col_index in range(cols):
            table.columns[col_index].width = Inches(2.2)

        image_y_offset = top + height + Inches(0.3)
        image_height = Inches(1.3)
        image_gap = Inches(0.2)

        for r in range(rows):
            for c in range(cols):
                cell = table.cell(r, c)
                cell.text = str(table_data[r][c])
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(10)
                if r == 0:
                    para.font.bold = True
                    para.font.color.rgb = RGBColor(255, 255, 255)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0, 112, 192)
                else:
                    para.font.color.rgb = RGBColor(0, 0, 0)

        # 🖼 Add product images below table
        for i in range(1, rows):  # Skip header row
            product_name = table_data[i][0]
            image_url = search_image_url_serpapi(product_name)
            if image_url:
                image_data = download_image(image_url)
                if image_data:
                    try:
                        slide.shapes.add_picture(image_data, Inches(0.5), image_y_offset, height=image_height)
                        print(f"🖼️ Inserted image for: {product_name}")
                        image_y_offset += image_height + image_gap
                    except Exception as e:
                        print(f"[Image Insert Error] {product_name}: {e}")
                else:
                    print(f"⚠️ Could not download image for {product_name}")
            else:
                print(f"⚠️ No image URL found for {product_name}")

        # 💰 Total price
        total = sum(p * q for p, q in price_qty_list)
        print(f"💰 Total Price for slide: ¥{total:,.0f}")

        # Estimate slide bottom margin
        slide_bottom_limit = Inches(7.2)  # Safe vertical limit to not overflow

        # Decide where to place price box
        price_top = image_y_offset + Inches(0.3)
        if price_top + Inches(0.5) > slide_bottom_limit:
            price_top = top + height + Inches(0.2)  # move it directly under table

        txBox = slide.shapes.add_textbox(Inches(6.5), price_top, Inches(2.5), Inches(0.8))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = f"Total: ¥{total:,.0f}"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 112, 192)
        p.alignment = PP_ALIGN.RIGHT

        add_logo(slide)
        add_footer(slide, index)
        return slide



    slide_index = 1

    # Slide 1: Title
    is_japanese = contains_japanese(quotation_text)

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    remove_title_placeholder(slide)
    set_slide_bg_color(slide)
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    tf = box.text_frame
    p = tf.add_paragraph()
    if is_japanese:
        p.text = "ハードウェア構成の見積もり" 
    else:
        p.text = "Hardware Configuration Quotations" 
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    subtitle = tf.add_paragraph()
    if is_japanese:
        subtitle.text = "大塚商会 AI営業エージェントによって生成されました"
    else:
        subtitle.text = "Generated by Otsuka Shokai AI Sales Agent"
    subtitle.font.size = Pt(18)
    subtitle.font.color.rgb = RGBColor(100, 100, 100)

    add_logo(slide)
    add_footer(slide, slide_index)
    slide_index += 1

    # Slide 2: Client Info
    # Slide 2: Client Info
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_bg_color(slide)


    slide.shapes.title.text = "クライアント情報" if is_japanese else "Client Information"
    underline_title(slide)

    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8), Inches(4))
    tf = box.text_frame
    tf.word_wrap = True

    if is_japanese:
        labels = [("name", "1. クライアント名"), ("company", "2. 会社名"), ("email", "3. メール"), ("contact", "4. 電話番号")]
    else:
        labels = [("name", "1. Client Name"), ("company", "2. Company"), ("email", "3. Contact Email"), ("contact", "4. Contact Number")]

    for key, label in labels:
        p = tf.add_paragraph()
        p.text = f"{label}: {client_info.get(key, '')}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(50, 50, 50)

        add_logo(slide)
        add_footer(slide, slide_index)
        slide_index += 1

    # --- Parse quotation_text ---
    lines = quotation_text.strip().splitlines()
    current_quotation = ""
    table_data = []
    price_qty_list = []
    recommendation_lines = []
    inside_quote = False

    for line in lines:
        line = line.strip()
        if line.startswith("## Quotation") or line.startswith("## 見積もり"):
            if table_data:
                add_quotation_slide(prs, current_quotation, table_data, price_qty_list, slide_index)
                slide_index += 1
            current_quotation = line.replace("##", "").strip()
            table_data = [["商品名", "仕様", "価格", "数量"]] if "見積もり" in line else [["Product Name", "Specs", "Price", "Qty"]]
            price_qty_list = []
            inside_quote = True

        elif line.startswith("Product Name:") or line.startswith("商品名:"):
            pname = line.split(":", 1)[1].strip()

        elif line.startswith("Specs:") or line.startswith("仕様:"):
            specs = line.split(":", 1)[1].strip()

        elif line.startswith("Price:") or line.startswith("価格:"):
            raw = line.split(":", 1)[1].strip()
            clean = re.sub(r"[^\d.]", "", raw)
            try:
                price = float(clean)
            except ValueError:
                price = 0.0

        elif line.startswith("Quantity:") or line.startswith("数量:"):
            qty_raw = line.split(":", 1)[1].strip()
            digits = re.findall(r'\d+', qty_raw)
            qty = int(digits[0]) if digits else 1
            table_data.append([pname, specs, f"¥{price:,.0f}", str(qty)])
            price_qty_list.append((price, qty))

        elif line.startswith("## Recommendation") or line.startswith("## 推奨案"):
            if table_data:
                add_quotation_slide(prs, current_quotation, table_data, price_qty_list, slide_index)
                slide_index += 1
            table_data = []
            inside_quote = False

        elif not inside_quote and line:
            recommendation_lines.append(line)

    if table_data and inside_quote:
        add_quotation_slide(prs, current_quotation, table_data, price_qty_list, slide_index)
        slide_index += 1

    # --- Recommendation Slide ---
    if recommendation_lines:
        chunks = [recommendation_lines[i:i + 10] for i in range(0, len(recommendation_lines), 10)]
        for idx, chunk in enumerate(chunks):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            set_slide_bg_color(slide)
            title = "ベスト推奨案" if any(re.search(r'[\u3040-\u30FF\u4E00-\u9FFF]', line) for line in chunk) else "Best Recommendation"
            if len(chunks) > 1:
                title += f" (Part {idx + 1})"
            slide.shapes.title.text = title
            underline_title(slide)

            box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(8.5), Inches(4.5))
            tf = box.text_frame
            tf.word_wrap = True
            font_size = Pt(20 if len(chunk) <= 6 else 16 if len(chunk) <= 10 else 14)
            for line in chunk:
                bullet = tf.add_paragraph()
                bullet.text = f"• {line}"
                bullet.level = 0
                bullet.font.size = font_size
                bullet.font.color.rgb = RGBColor(60, 60, 60)

            add_logo(slide)
            add_footer(slide, slide_index)
            slide_index += 1

    # --- Final Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_bg_color(slide)
    remove_title_placeholder(slide)

    box = slide.shapes.add_textbox(Inches(2), Inches(1.8), Inches(6), Inches(1))
    p = box.text_frame.paragraphs[0]
    if is_japanese:
        p.text = "----ありがとうございます----"
    else:
        p.text = "----Thank You----"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.underline = True
    p.alignment = PP_ALIGN.CENTER

    info_box = slide.shapes.add_textbox(Inches(1.5), Inches(3), Inches(7), Inches(2.5))
    tf = info_box.text_frame
    if is_japanese:
        thank_you_lines = (
        [
            "- Otsuka Shokai にご関心をお寄せいただきありがとうございます。",
            "- ご質問がございましたら、以下までご連絡ください：",
            "📧 support@otsuka-shokai.co.jp",
            "🌐 www.otsuka-shokai.co.jp"
        ])
    else:
        thank_you_lines=(
            [
            "We appreciate your interest in Otsuka Shokai.",
            "For any inquiries, reach out at:",
            "📧 support@otsuka-shokai.co.jp",
            "🌐 www.otsuka-shokai.co.jp"
        ])
    for line in thank_you_lines:
        para = tf.add_paragraph()
        para.text = line
        para.font.size = Pt(18)
        para.alignment = PP_ALIGN.CENTER
        
    add_logo(slide)
    add_footer(slide, slide_index)

    pptx_buffer = BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    print(f"PPT time: {time.time() - starttimeppt:.2f} sec")
    return pptx_buffer
