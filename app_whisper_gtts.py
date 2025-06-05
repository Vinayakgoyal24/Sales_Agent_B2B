import streamlit as st
import os
import pandas as pd
from dotenv import load_dotenv
from langchain_community.document_loaders import TextLoader
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_chroma import Chroma
from langchain_openai import AzureOpenAIEmbeddings, AzureChatOpenAI
from langchain.prompts import ChatPromptTemplate
from langgraph.graph import START, StateGraph
from typing_extensions import List, TypedDict
import tiktoken
from io import BytesIO
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER
from reportlab.pdfbase.pdfmetrics import stringWidth
from datetime import datetime, timedelta
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from io import BytesIO
from streamlit_mic_recorder import mic_recorder        # ← NEW
import tempfile                                        # ← NEW
import whisper                                         # ← NEW
from gtts import gTTS                                  # ← NEW
import torch

import shutil, os, streamlit as st
st.sidebar.write("ffmpeg on PATH →", shutil.which("ffmpeg"))
st.sidebar.write("first 500 chars of PATH →", os.environ["PATH"][:500])

# Load env vars
load_dotenv()

# --- Initialize Azure Clients ---
llm = AzureChatOpenAI(
    azure_endpoint=os.getenv("AZURE_OPENAI_ENDPOINT"),
    azure_deployment=os.getenv("AZURE_OPENAI_DEPLOYMENT_NAME"),
    openai_api_version=os.getenv("AZURE_OPENAI_API_VERSION"),
)

embeddings = AzureOpenAIEmbeddings(
    azure_endpoint=os.getenv("AZURE_OPENAI_EMBEDDING_ENDPOINT"),
    azure_deployment=os.getenv("AZURE_OPENAI_EMBEDDING_DEPLOYMENT_NAME"),
    openai_api_version=os.getenv("AZURE_OPENAI_EMBEDDING_API_VERSION"),
)

vector_store = Chroma(
    collection_name="example_collection",
    embedding_function=embeddings,
    persist_directory="./chroma_langchain_db",
)

# --- Load CSVs and Create Documents ---
def load_csv_as_documents(folder_path="data"):
    documents = []
    for file_name in os.listdir(folder_path):
        if file_name.endswith(".csv"):
            df = pd.read_csv(os.path.join(folder_path, file_name)).head(20)
            for _, row in df.iterrows():
                content = "\n".join([f"{k}: {v}" for k, v in row.items()])
                documents.append(Document(page_content=content, metadata={"source": file_name}))
    return documents

if vector_store._collection.count() == 0:
    with st.spinner("Indexing documents..."):
        docs = load_csv_as_documents("data")
        splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        splits = splitter.split_documents(docs)
        vector_store.add_documents(splits)
        vector_store.persist()


# --- Initialize Whisper Model ---
@st.cache_resource(show_spinner="Loading Whisper model…")
def load_whisper_local(model_name: str = "base"):
    return whisper.load_model(
        model_name,
        device="cuda" if torch.cuda.is_available() else "cpu"
    )

def transcribe_audio(wav_input) -> str:
    """
    Accepts either the dict returned by mic_recorder or raw bytes.
    Returns the Whisper transcription string.
    """
    # ── 1. normalise to raw bytes ──────────────────────────────────────────
    if isinstance(wav_input, dict):            # mic_recorder style
        wav_bytes = wav_input.get("bytes", b"")
    else:                                      # already bytes
        wav_bytes = wav_input

    if not wav_bytes:
        return ""                              # nothing to transcribe

    # ── 2. write to temp file and run Whisper ──────────────────────────────
    model = load_whisper_local()
    with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as tmp:
        tmp.write(wav_bytes)
        tmp_path = tmp.name
    try:
        result = model.transcribe(tmp_path)    # {'text': '...'}
        return result["text"].strip()
    finally:
        os.remove(tmp_path)


def _preprocess_for_tts(raw: str) -> str:
    """
    Prepare text for gTTS:
    • Remove leading markdown symbols (##, bullets) but keep the wording.
    • Collapse multiple lines into full sentences separated by a pause.
    """
    cleaned_lines = []
    for line in raw.splitlines():
        line = line.strip()

        # Keep heading wording but drop the hashes, e.g. "## Quotation 1"
        if line.startswith("##"):
            line = line.lstrip("#").strip()

        # Remove common bullet prefixes
        line = re.sub(r"^[\-\*\•]\s*", "", line)

        if line:                       # skip empty lines
            cleaned_lines.append(line)

    # Join with a double space so gTTS inserts a brief pause
    return ".  ".join(cleaned_lines)



def tts_stream(text: str, lang: str = "en", tld: str = "co.uk") -> BytesIO:
    """
    Convert *text* to speech with gTTS and return an MP3 BytesIO.
    • `tld` picks the Google voice accent ─ try 'com', 'co.uk', 'com.au', etc.
    """
    processed = _preprocess_for_tts(text)
    mp3 = BytesIO()
    gTTS(text=processed, lang=lang, tld=tld, slow=False).write_to_fp(mp3)
    mp3.seek(0)
    return mp3


# --- Prompt Template ---
prompt = ChatPromptTemplate.from_messages([
    ("system",
     "You are a professional hardware sales assistant at Otsuka Corporation. Based on the user's request and the context, provide 2–3 detailed hardware configuration quotations. "
     "Each quotation should include:\n"
     "- Product Name\n- Specs\n- Price\n- Quantity\n- Total Price\n\n"
     "Use this structure:\n"
     "## Quotation 1\nProduct Name: ...\nSpecs: ...\nPrice: ...\nQuantity: ...\n...\nTotal Price: ...\n"
     "## Quotation 2 ...\n\n"
     "Then provide a clear comparison of the quotations and recommend the best one based on:\n"
     "- Price\n- Suitability for the user's need\n- Performance vs cost.\n"
     "Use a section titled:\n"
     "## Recommendation\n"
     "Mention why the chosen quote is the best and highlight key differences with others.\n\n"
     "Keep tone professional and brief. Do not fabricate information if context is insufficient."),
    ("human", "Question: {question}\n\nContext:\n{context}")
])


from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from PIL import Image as PILImage
import io


def get_slide_preview(pptx_bytes: BytesIO) -> BytesIO:
    from comtypes import client
    import tempfile

    # Save .pptx to temp file
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        tmp.write(pptx_bytes.read())
        ppt_path = tmp.name

    # Export first slide to PNG (only works on Windows with PowerPoint installed)
    powerpoint = client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1
    deck = powerpoint.Presentations.Open(ppt_path)
    tmp_img = ppt_path.replace(".pptx", "_1.png")
    deck.SaveAs(tmp_img, 18)  # 18 for PNG
    deck.Close()
    powerpoint.Quit()

    # Load PNG preview into BytesIO
    with open(tmp_img, "rb") as img_file:
        return BytesIO(img_file.read())


def generate_slides(quotation_text: str) -> BytesIO:
    prs = Presentation()
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    def add_logo(slide):
        if os.path.exists("otsuka_im.png"):
            LOGO_WIDTH_INCHES = 1.2  # Adjust width
            LOGO_TOP_MARGIN = 0.2
            LOGO_RIGHT_MARGIN = 0.3

        # Calculate position from right edge
            left = slide_width - Inches(LOGO_RIGHT_MARGIN + LOGO_WIDTH_INCHES)
            top = Inches(LOGO_TOP_MARGIN)

            slide.shapes.add_picture("otsuka_im.png", left, top, width=Inches(LOGO_WIDTH_INCHES))


    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Hardware Configuration Quotations"
    slide.placeholders[1].text = "Generated by Otsuka Corporation's Sales Assistant"
    add_logo(slide)

    # Client Info Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "📋 Client Information"
    content = slide.placeholders[1]
    content.text = "\n".join([
        "Client Name: Acme Solutions Pvt. Ltd.",
        "Client Address: 123, Innovation Tower, Marunouchi, Tokyo",
        "Contact Person: John Doe",
        "Email: john.doe@example.com",
        "Phone: +81 90-1234-5678"
    ])
    add_logo(slide)

    lines = quotation_text.strip().splitlines()
    current_quotation = ""
    recommendation_lines = []
    table_data = []
    price_qty_list = []
    inside_quote = False

    for line in lines:
        line = line.strip()

        if line.startswith("## Quotation"):
            if table_data:
                add_quotation_slide(prs, current_quotation, table_data, price_qty_list)
            current_quotation = line.replace("##", "").strip()
            table_data = [["Product Name", "Specs", "Price", "Qty"]]
            price_qty_list = []
            inside_quote = True

        elif line.startswith("Product Name:"):
            pname = line.split(":", 1)[1].strip()

        elif line.startswith("Specs:"):
            specs = line.split(":", 1)[1].strip()

        elif line.startswith("Price:"):
            raw = line.split(":", 1)[1].strip()
            clean = re.sub(r"[^\d\.]", "", raw)
            price = float(clean)

        elif line.startswith("Quantity:"):
            qty = int(line.split(":", 1)[1].strip())
            table_data.append([pname, specs, f"${price:,.0f}", str(qty)])
            price_qty_list.append((price, qty))

        elif line.startswith("## Recommendation"):
            if table_data:
                add_quotation_slide(prs, current_quotation, table_data, price_qty_list)
                table_data = []
            inside_quote = False

        elif not inside_quote and line:
            recommendation_lines.append(line)

    if table_data and inside_quote:
        add_quotation_slide(prs, current_quotation, table_data, price_qty_list)

    # Recommendation Slide
    if recommendation_lines:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "🎯 Recommendation"
        content = slide.placeholders[1]
        content.text = "\n".join(recommendation_lines)
        add_logo(slide)

    # Thank You Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "🙏 Thank You"
    content = slide.placeholders[1]
    content.text = "\n".join([
        "Otsuka Corporation",
        "Head Office: 2-18-4 Iidabashi, Chiyoda-ku, Tokyo 102-8573",
        "Website: https://www.otsuka-shokai.co.jp"
    ])
    add_logo(slide)

    # Save to buffer
    pptx_buffer = BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    return pptx_buffer

def add_quotation_slide(prs, title, table_data, price_qty_list):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    shapes = slide.shapes
    shapes.title.text = title

    rows = len(table_data)
    cols = len(table_data[0])
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(0.8 + 0.4 * rows)

    table = shapes.add_table(rows, cols, left, top, width, height).table

    for col_index in range(cols):
        table.columns[col_index].width = Inches(2.2)

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(table_data[r][c])
            cell.text_frame.paragraphs[0].font.size = Pt(11)

    total = sum(p * q for p, q in price_qty_list)
    txBox = slide.shapes.add_textbox(Inches(0.5), top + height + Inches(0.3), Inches(5), Inches(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = f"Total Price: ¥{total:,.0f}"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)
def generate_pdf(quotation_text: str) -> BytesIO:
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer,
                            pagesize=A4,
                            rightMargin=30, leftMargin=30,
                            topMargin=30, bottomMargin=30)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(name="TitleCenter", parent=styles["Title"], alignment=TA_CENTER)
    highlight_style = ParagraphStyle(name="Highlight", parent=styles["Normal"], textColor=colors.green, fontSize=12)

    elements = []

    # --- Header: Company Info ---
    logo_path = "otsuka_im.png"
    if os.path.exists(logo_path):
        logo = Image(logo_path, width=100, height=50)
        logo.hAlign = "LEFT"
        elements.append(logo)

    elements.append(Paragraph("Otsuka Corporation", title_style))
    elements.append(Paragraph("Head Office, 2-18-4 Iidabashi, Chiyoda-ku, Tokyo 102-8573", styles["Normal"]))
    elements.append(Paragraph("Website: <a href='https://www.otsuka-shokai.co.jp/'>www.otsuka-shokai.co.jp</a>",
                              styles["Normal"]))
    elements.append(Spacer(1, 12))

    # --- Quotation Info ---
    elements.append(Paragraph("<b>🧾 Quotation</b>", styles["Heading2"]))
    today = datetime.now()
    validity = today + timedelta(days=7)
    elements.append(Paragraph(f"Date of Issue: {today.strftime('%Y-%m-%d')}", styles["Normal"]))
    elements.append(Paragraph(f"Validity: {validity.strftime('%Y-%m-%d')} (7 days)", styles["Normal"]))
    elements.append(Spacer(1, 12))

    # --- Client Info (Dummy) ---
    elements.append(Paragraph("<b>📋 Client Information</b>", styles["Heading3"]))
    client_info = [
        "Client Name: Acme Solutions Pvt. Ltd.",
        "Client Address: 123, Innovation Tower, Marunouchi, Tokyo",
        "Contact Person: John Doe",
        "Email: john.doe@example.com",
        "Phone: +81 90-1234-5678"
    ]
    for line in client_info:
        elements.append(Paragraph(line, styles["Normal"]))
    elements.append(Spacer(1, 12))

    # --- Parse LLM output and build tables ---
    lines = quotation_text.strip().splitlines()
    table_data = []
    total_prices = []
    price_qty_list = []
    recommendation_lines = []
    current_quotation = ""
    inside_quote = False

    def build_table(title, data, bg_color):
        table_style = styles = getSampleStyleSheet()
        cell_style = styles["Normal"]
        font_name = "Helvetica"
        font_size = 10

        # Convert data cells to Paragraphs (except header)
        tbl = [data[0]]
        for row in data[1:]:
            tbl.append([Paragraph(str(cell), cell_style) for cell in row])

        # Dynamic column widths
        transposed = list(zip(*data))
        col_widths = []
        for col in transposed:
            max_w = max(stringWidth(str(item), font_name, font_size) for item in col)
            col_widths.append(min(max_w + 20, 200))

        table = Table(tbl, hAlign="LEFT", colWidths=col_widths)
        table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), bg_color),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.whitesmoke),
            ("GRID", (0, 0), (-1, -1), 0.7, colors.black),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), font_size),
            ("ALIGN", (2, 1), (2, -1), "RIGHT"),
            ("ALIGN", (3, 1), (3, -1), "CENTER"),
            ("VALIGN", (0, 1), (-1, -1), "TOP"),
        ]))

        elements.append(Paragraph(f"<b>{title}</b>", styles["Heading4"]))
        elements.append(table)

    for line in lines:
        line = line.strip()

        if line.startswith("## Quotation"):
            # flush previous quote
            if table_data:
                subtotal = sum(p * q for p, q in price_qty_list)
                subtotal_str = f"{subtotal:,.0f}"
                build_table(current_quotation, table_data, bg_color=colors.HexColor("#4472C4"))
                elements.append(Paragraph(f"<b>Total Price ({current_quotation}):</b> ¥{subtotal_str}", styles["Normal"]))
                total_prices.append((current_quotation, subtotal_str))
                elements.append(Spacer(1, 10))

            current_quotation = line.replace("##", "").strip()
            table_data = [["Product Name", "Specs", "Price ($)", "Qty"]]
            price_qty_list = []
            inside_quote = True

        elif line.startswith("Product Name:"):
            pname = line.split(":", 1)[1].strip()

        elif line.startswith("Specs:"):
            specs = line.split(":", 1)[1].strip()

        elif line.startswith("Price:"):
            raw = line.split(":", 1)[1].strip()
            # strip out anything but digits and dot
            clean = re.sub(r"[^\d\.]", "", raw)
            price = float(clean)


        elif line.startswith("Quantity:"):
            qty = int(line.split(":", 1)[1].strip())
            table_data.append([pname, specs, f"{price:,.0f}", str(qty)])
            price_qty_list.append((price, qty))

        elif line.startswith("## Recommendation"):
            # flush last quote before recommendation
            if table_data and inside_quote:
                subtotal = sum(p * q for p, q in price_qty_list)
                subtotal_str = f"{subtotal:,.0f}"
                build_table(current_quotation, table_data, bg_color=colors.HexColor("#4472C4"))
                elements.append(Paragraph(f"<b>Total Price ({current_quotation}):</b> ¥{subtotal_str}", styles["Normal"]))
                total_prices.append((current_quotation, subtotal_str))
                elements.append(Spacer(1, 10))
            inside_quote = False
            elements.append(Paragraph("<b>🎯 Recommendation</b>", styles["Heading3"]))

        elif not inside_quote and line:
            recommendation_lines.append(line)

    # Summary tables if any left unflushed
    if table_data and inside_quote:
        build_table(current_quotation, table_data, bg_color=colors.HexColor("#4472C4"))

    # Pricing Summary
    elements.append(Spacer(1, 12))
    elements.append(Paragraph("<b>📊 Pricing Summary</b>", styles["Heading3"]))
    for qname, tprice in total_prices:
        elements.append(Paragraph(f"• {qname}: ¥{tprice}", styles["Normal"]))

    # Recommendation Section
    if recommendation_lines:
        elements.append(Spacer(1, 12))
        elements.append(Paragraph("<b>✅ Best Recommendation</b>", styles["Heading3"]))
        for line in recommendation_lines:
            elements.append(Paragraph(line, highlight_style))

    doc.build(elements)
    buffer.seek(0)
    return buffer

# --- LangGraph App Logic ---
class State(TypedDict):
    question: str
    context: List[Document]
    answer: str

def retrieve(state: State):
    return {"context": vector_store.similarity_search(state["question"])}

def generate(state: State):
    docs_content = "\n\n".join(doc.page_content for doc in state["context"])
    messages = prompt.invoke({"question": state["question"], "context": docs_content})
    response = llm.invoke(messages)
    return {"answer": response.content}

graph_builder = StateGraph(State).add_sequence([retrieve, generate])
graph_builder.add_edge(START, "retrieve")
graph = graph_builder.compile()

# --- Streamlit UI ---
st.title("💻 Computer Hardware Sales Assistant")

if "result" not in st.session_state:
    st.session_state.result = None

if "text_query" not in st.session_state:
    st.session_state.text_query = ""


# user_query = st.text_input(
#     "Enter your query:",
#     placeholder="E.g., Best PC setup for video editing…",
#     key="text_query"
# )

# ── TEXT + MIC ROW ──────────────────────────────────────────────────────────
# Put the text box and the mic icon on the same line
# col_text, col_mic = st.columns([8, 1])

# with col_text:
#     # keep the shared key so both voice and typing update the same value
#     user_query = st.text_input(
#         "Enter your query:",
#         placeholder="E.g., Best PC setup for video editing…",
#         key="text_query"
#     )

# with col_mic:
#     # mic_recorder shows a round mic icon when start_prompt/stop_prompt = None
#     audio_bytes = mic_recorder(
#         start_prompt=None,      # default mic icon
#         stop_prompt=None,       # default stop icon
#         key="recorder_button"
#     )

# # If the user just finished talking, transcribe → populate the text box
# if audio_bytes:
#     with st.spinner("Transcribing…"):
#         spoken_text = transcribe_audio(audio_bytes)
#         if spoken_text:
#             st.success(f"You said: {spoken_text}")
#             # write into the text box and refresh the page so it appears there
#             st.session_state.text_query = spoken_text
#             st.experimental_rerun()


# if user_query and st.button("💬 Get Recommendation"):
#     with st.spinner("Processing your query..."):
#         st.session_state.result = graph.invoke({"question": user_query})

# if st.session_state.result:
#     st.subheader("💡 Suggested Answer")
#     st.write(st.session_state.result["answer"])

#     enc = tiktoken.encoding_for_model("gpt-4")
#     input_tokens = len(enc.encode(user_query))
#     output_tokens = len(enc.encode(st.session_state.result["answer"]))
#     st.markdown(f"🔢 Input tokens: {input_tokens} | Output tokens: {output_tokens} | Total: {input_tokens + output_tokens}")

#     pdf_bytes = generate_pdf(st.session_state.result["answer"])
#     st.download_button(
#         label="📄 Download Quotation as PDF",
#         data=pdf_bytes,
#         file_name="hardware_quotation.pdf",
#         mime="application/pdf"
#     )



#     slide_bytes = generate_slides(st.session_state.result["answer"])
#     st.download_button(
#         label="📊 Download Slides (PPTX)",
#         data=slide_bytes,
#         file_name="hardware_quotation.pptx",
#         mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
#     )

#     # OPTIONAL: Show slide preview (first slide)
#     try:
#         preview_image = get_slide_preview(BytesIO(slide_bytes.read()))
#         st.image(preview_image, caption="📽️ Slide Preview", use_column_width=True)
#     except Exception as e:
#         st.warning("Could not render slide preview. (Only works on Windows with PowerPoint installed.)")

#     if st.button("🔊 Speak this quotation"):
#         with st.spinner("Generating speech…"):
#             mp3 = tts_stream(st.session_state.result["answer"])
#             st.audio(mp3.read(), format="audio/mp3")


# ───────────────────────────────────────────────────────────────────────────────
# MIC + TRANSCRIPTION + TEXT INPUT (replace your existing block with this)
# ───────────────────────────────────────────────────────────────────────────────

# 1) Ensure the flag exists
if "just_transcribed" not in st.session_state:
    st.session_state.just_transcribed = False

# 2) Mic recorder icon
audio_bytes = mic_recorder(
    start_prompt="🎤",   # show a mic emoji
    stop_prompt="⏹️",
    key="recorder_button"
)

# 3) Only transcribe once per “spoken” event
if audio_bytes and not st.session_state.just_transcribed:
    with st.spinner("Transcribing…"):
        spoken_text = transcribe_audio(audio_bytes)
        if spoken_text:
            st.success(f"You said: {spoken_text}")
            st.session_state.text_query = spoken_text
            st.session_state.just_transcribed = True

# 4) Define a callback to allow manual edits to reset the flag
def reset_transcription_flag():
    st.session_state.just_transcribed = False

# 5) The text_input uses the same key "text_query" and resets the flag on edit
user_query = st.text_input(
    "Enter your query:",
    placeholder="E.g., Best PC setup for video editing…",
    key="text_query",
    on_change=reset_transcription_flag
)
# ───────────────────────────────────────────────────────────────────────────────


# 3) “Get Recommendation” button
if user_query and st.button("💬 Get Recommendation"):
    with st.spinner("Processing your query..."):
        st.session_state.result = graph.invoke({"question": user_query})

# 4) Show result + downloads + TTS
if st.session_state.result:
    answer = st.session_state.result["answer"]
    st.subheader("💡 Suggested Answer")
    st.write(answer)

    enc = tiktoken.encoding_for_model("gpt-4")
    input_tokens = len(enc.encode(user_query))
    output_tokens = len(enc.encode(answer))
    st.markdown(f"🔢 Input tokens: {input_tokens} | Output tokens: {output_tokens} | Total: {input_tokens + output_tokens}")

    pdf_bytes = generate_pdf(answer)
    st.download_button(
        label="📄 Download Quotation as PDF",
        data=pdf_bytes,
        file_name="hardware_quotation.pdf",
        mime="application/pdf"
    )

    slide_bytes = generate_slides(answer)
    st.download_button(
        label="📊 Download Slides (PPTX)",
        data=slide_bytes,
        file_name="hardware_quotation.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

    try:
        preview_image = get_slide_preview(BytesIO(slide_bytes.read()))
        st.image(preview_image, caption="📽️ Slide Preview", use_column_width=True)
    except Exception:
        pass

    if st.button("🔊 Speak this quotation"):
        with st.spinner("Generating speech…"):
            mp3 = tts_stream(answer)
            st.audio(mp3.read(), format="audio/mp3")

