import os
import re
import pandas as pd
import tiktoken
from dotenv import load_dotenv
from io import BytesIO
import streamlit as st
from langchain_community.document_loaders import TextLoader
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_chroma import Chroma
from langchain_openai import AzureOpenAIEmbeddings, AzureChatOpenAI
from langchain.prompts import ChatPromptTemplate
from langgraph.graph import START, StateGraph
from typing_extensions import List, TypedDict
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER
from reportlab.pdfbase.pdfmetrics import stringWidth
from datetime import datetime, timedelta
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from email.mime.base import MIMEBase
from email import encoders


# Load env vars
load_dotenv()

# --- Initialize Azure Clients ---
llm = AzureChatOpenAI(
    azure_endpoint=os.getenv("AZURE_OPENAI_ENDPOINT"),
    azure_deployment=os.getenv("AZURE_OPENAI_DEPLOYMENT_NAME"),
    openai_api_version=os.getenv("AZURE_OPENAI_API_VERSION"),
)

embeddings = AzureOpenAIEmbeddings(
    azure_endpoint=os.getenv("AZURE_OPENAI_EMBEDDING_ENDPOINT"),
    azure_deployment=os.getenv("AZURE_OPENAI_EMBEDDING_DEPLOYMENT_NAME"),
    openai_api_version=os.getenv("AZURE_OPENAI_EMBEDDING_API_VERSION"),
)

vector_store = Chroma(
    collection_name="example_collection",
    embedding_function=embeddings,
    persist_directory="./chroma_langchain_db",
)

# --- Load CSVs and Create Documents ---
def load_csv_as_documents(folder_path="data"):
    documents = []
    for file_name in os.listdir(folder_path):
        if file_name.endswith(".csv"):
            df = pd.read_csv(os.path.join(folder_path, file_name)).head(20)
            for _, row in df.iterrows():
                content = "\n".join([f"{k}: {v}" for k, v in row.items()])
                documents.append(Document(page_content=content, metadata={"source": file_name}))
    return documents

if vector_store._collection.count() == 0:
    with st.spinner("Indexing documents..."):
        docs = load_csv_as_documents("data")
        splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        splits = splitter.split_documents(docs)
        vector_store.add_documents(splits)
        vector_store.persist()

# --- Prompt Template ---
prompt = ChatPromptTemplate.from_messages([
    ("system",
     "You are a professional hardware sales assistant at Otsuka Corporation. Based on the user's request and the context, provide 2–3 detailed hardware configuration quotations. "
     "Each quotation should include:\n"
     "- Product Name\n- Specs\n- Price\n- Quantity\n- Total Price\n\n"
     "Use this structure:\n"
     "## Quotation 1\nProduct Name: ...\nSpecs: ...\nPrice: ...\nQuantity: ...\n...\nTotal Price: ...\n"
     "## Quotation 2 ...\n\n"
     "Then provide a clear comparison of the quotations and recommend the best one based on:\n"
     "- Price\n- Suitability for the user's need\n- Performance vs cost.\n"
     "Use a section titled:\n"
     "## Recommendation\n"
     "Mention why the chosen quote is the best and highlight key differences with others.\n\n"
     "Keep tone professional and brief. Do not fabricate information if context is insufficient."),
    ("human", "Question: {question}\n\nContext:\n{context}")
])


from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from PIL import Image as PILImage
import io


def send_email_with_attachment(to_email: str, subject: str, body: str):

    smtp_port = 587
    smtp_server = "smtp.gmail.com"
    sender_email = "vinayak.otsuka@gmail.com"
    pswd = "djjvyfubleftjmwh"

    if not sender_email or not pswd:
        return False, "Missing sender email or password in environment variables."

    msg = MIMEMultipart()
    msg['From'] = sender_email
    msg['To'] = to_email
    msg['Subject'] = subject

    msg.attach(MIMEText(body, 'plain'))

    filename = "static/hardware_quotation.pdf"
    try:
        with open(filename, 'rb') as attachment:
            attachment_package = MIMEBase('application', 'octet-stream')
            attachment_package.set_payload(attachment.read())
            encoders.encode_base64(attachment_package)
            attachment_package.add_header('Content-Disposition', f"attachment; filename={filename}")
            msg.attach(attachment_package)
    except FileNotFoundError:
        return False, f"Attachment file not found: {filename}"

    text = msg.as_string()

    try:
        print("Connecting to server...")
        TIE_server = smtplib.SMTP(smtp_server, smtp_port)
        TIE_server.starttls()
        TIE_server.login(sender_email, pswd)
        print("Successfully connected to server")

        TIE_server.sendmail(sender_email, to_email, text)
        TIE_server.quit()
        return True, f"Email successfully sent to {to_email}"
    except Exception as e:
        return False, f"Failed to send email to {to_email}. Error: {e}"


from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx import Presentation
from io import BytesIO
import os
import re
from pptx.enum.shapes import MSO_SHAPE

def generate_slides(quotation_text: str, client_info: dict) -> BytesIO:
    prs = Presentation()
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    bg_color = RGBColor(0xFF, 0xED, 0xEC)  # Soft pink

    def set_slide_bg_color(slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    def add_logo(slide):
        if os.path.exists("otsuka_im.png"):
            LOGO_WIDTH_INCHES = 1.2
            LOGO_TOP_MARGIN = 0.2
            LOGO_RIGHT_MARGIN = 0.3
            left = slide_width - Inches(LOGO_RIGHT_MARGIN + LOGO_WIDTH_INCHES)
            top = Inches(LOGO_TOP_MARGIN)
            slide.shapes.add_picture("otsuka_im.png", left, top, width=Inches(LOGO_WIDTH_INCHES))
    
    def add_quotation_slide(prs, title, table_data, price_qty_list):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        set_slide_bg_color(slide)
        shapes = slide.shapes
        shapes.title.text = title

        rows = len(table_data)
        cols = len(table_data[0])
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(0.8 + 0.4 * rows)

        table = shapes.add_table(rows, cols, left, top, width, height).table

        for col_index in range(cols):
            table.columns[col_index].width = Inches(2.2)

        for r in range(rows):
            for c in range(cols):
                cell = table.cell(r, c)
                cell.text = str(table_data[r][c])
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(11)

                if r == 0:  # Header row
                    para.font.bold = True
                    para.font.color.rgb = RGBColor(0, 0, 0)  # Black

    # Total Price Text
        total = sum(p * q for p, q in price_qty_list)
        txBox = slide.shapes.add_textbox(Inches(0.5), top + height + Inches(0.3), Inches(5), Inches(1))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = f"- Total Price: ¥{total:,.0f}"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 112, 192)  # Blue

        add_logo(slide)

    # Slide: Title
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_bg_color(slide)
    slide.shapes.title.text = "Hardware Configuration Quotations"
    add_logo(slide)

    # Slide: Client Information
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_bg_color(slide)
    title = slide.shapes.title
    title.text = "--- Client Information ---"
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8.5), Inches(4))
    tf = content_box.text_frame
    tf.word_wrap = True
    for key, label in [("name", "1. Client Name"), ("2. Company", "Client Company"), ("3. Contact_Email", "Email"), ("4. Contact_Phone", "Phone")]:
        p = tf.add_paragraph()
        p.text = f"{label}: {client_info.get(key, '')}"
        p.font.size = Pt(18)
    add_logo(slide)

    # Slide Parsing Logic
    lines = quotation_text.strip().splitlines()
    current_quotation = ""
    recommendation_lines = []
    table_data = []
    price_qty_list = []
    inside_quote = False

    for line in lines:
        line = line.strip()

        if line.startswith("## Quotation"):
            if table_data:
                add_quotation_slide(prs, current_quotation, table_data, price_qty_list)
            current_quotation = line.replace("##", "").strip()
            table_data = [["Product Name", "Specs", "Price", "Qty"]]
            price_qty_list = []
            inside_quote = True

        elif line.startswith("Product Name:"):
            pname = line.split(":", 1)[1].strip()

        elif line.startswith("Specs:"):
            specs = line.split(":", 1)[1].strip()

        elif line.startswith("Price:"):
            raw = line.split(":", 1)[1].strip()
            clean = re.sub(r"[^\d\.]", "", raw)
            price = float(clean)

        elif line.startswith("Quantity:"):
            qty = int(line.split(":", 1)[1].strip())
            table_data.append([pname, specs, f"${price:,.0f}", str(qty)])
            price_qty_list.append((price, qty))

        elif line.startswith("## Recommendation"):
            if table_data:
                add_quotation_slide(prs, current_quotation, table_data, price_qty_list)
                table_data = []
            inside_quote = False

        elif not inside_quote and line:
            recommendation_lines.append(line)

    if table_data and inside_quote:
        add_quotation_slide(prs, current_quotation, table_data, price_qty_list)

    # Slide: Recommendation
    if recommendation_lines:
        MAX_LINES_PER_SLIDE = 10
        chunks = [recommendation_lines[i:i + MAX_LINES_PER_SLIDE] for i in range(0, len(recommendation_lines), MAX_LINES_PER_SLIDE)]

        for chunk_index, chunk in enumerate(chunks):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            set_slide_bg_color(slide)
            title_text = "Best Recommendation"
            if len(chunks) > 1:
                title_text += f" (Part {chunk_index + 1})"
            slide.shapes.title.text = title_text

            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8.5), Inches(4.5))
            tf = content_box.text_frame
            tf.word_wrap = True

            # Dynamically scale font size
            font_size = Pt(20 if len(chunk) <= 6 else 16 if len(chunk) <= 10 else 14)

            for line in chunk:
                bullet = tf.add_paragraph()
                bullet.text = line
                bullet.level = 0
                bullet.font.size = font_size

            add_logo(slide)

    # Slide: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_bg_color(slide)
    slide.shapes.title.text = "--- Thank You ---"
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8.5), Inches(4.5))
    tf = content_box.text_frame
    for line in [
        "1. Otsuka Corporation",
        "2. Head Office: 2-18-4 Iidabashi, Chiyoda-ku, Tokyo 102-8573",
        "3. Website: https://www.otsuka-shokai.co.jp"
    ]:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
    add_logo(slide)

    pptx_buffer = BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    return pptx_buffer

def generate_pdf(quotation_text: str, client_info: dict) -> BytesIO:
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer,
                            pagesize=A4,
                            rightMargin=30, leftMargin=30,
                            topMargin=30, bottomMargin=30)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(name="TitleCenter", parent=styles["Title"], alignment=TA_CENTER)
    highlight_style = ParagraphStyle(name="Highlight", parent=styles["Normal"], textColor=colors.green, fontSize=12)

    elements = []

    # --- Header: Company Info ---
    logo_path = "otsuka_im.png"
    if os.path.exists(logo_path):
        logo = Image(logo_path, width=100, height=50)
        logo.hAlign = "LEFT"
        elements.append(logo)

    elements.append(Paragraph("Otsuka Corporation", title_style))
    elements.append(Paragraph("Head Office, 2-18-4 Iidabashi, Chiyoda-ku, Tokyo 102-8573", styles["Normal"]))
    elements.append(Paragraph("Website: <a href='https://www.otsuka-shokai.co.jp/'>www.otsuka-shokai.co.jp</a>",
                              styles["Normal"]))
    elements.append(Spacer(1, 12))

    # --- Quotation Info ---
    elements.append(Paragraph("<b>--- Quotation ---</b>", styles["Heading2"]))
    today = datetime.now()
    validity = today + timedelta(days=7)
    elements.append(Paragraph(f"Date of Issue: {today.strftime('%Y-%m-%d')}", styles["Normal"]))
    elements.append(Paragraph(f"Validity: {validity.strftime('%Y-%m-%d')} (7 days)", styles["Normal"]))
    elements.append(Spacer(1, 12))

    # --- Client Info (Dummy) ---
    elements.append(Paragraph("<b>--- Client Information ---</b>", styles["Heading3"]))
    client_info = [
        "1. Client Name: "+ client_info.get("name",""),
        "2. Client Company:"+ client_info.get("company",""),
        "3. Contact Person: John Doe",
        "4. Email: "+ client_info.get("email",""),
        "5. Phone: "+ client_info.get("phone","")
    ]
    for line in client_info:
        elements.append(Paragraph(line, styles["Normal"]))
    elements.append(Spacer(1, 12))

    # --- Parse LLM output and build tables ---
    lines = quotation_text.strip().splitlines()
    table_data = []
    total_prices = []
    price_qty_list = []
    recommendation_lines = []
    current_quotation = ""
    inside_quote = False

    def build_table(title, data, bg_color):
        table_style = styles = getSampleStyleSheet()
        cell_style = styles["Normal"]
        font_name = "Helvetica"
        font_size = 10

        # Convert data cells to Paragraphs (except header)
        tbl = [data[0]]
        for row in data[1:]:
            tbl.append([Paragraph(str(cell), cell_style) for cell in row])

        # Dynamic column widths
        transposed = list(zip(*data))
        col_widths = []
        for col in transposed:
            max_w = max(stringWidth(str(item), font_name, font_size) for item in col)
            col_widths.append(min(max_w + 20, 200))

        table = Table(tbl, hAlign="LEFT", colWidths=col_widths)
        table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), bg_color),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.whitesmoke),
            ("GRID", (0, 0), (-1, -1), 0.7, colors.black),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), font_size),
            ("ALIGN", (2, 1), (2, -1), "RIGHT"),
            ("ALIGN", (3, 1), (3, -1), "CENTER"),
            ("VALIGN", (0, 1), (-1, -1), "TOP"),
        ]))

        elements.append(Paragraph(f"<b>{title}</b>", styles["Heading4"]))
        elements.append(table)

    for line in lines:
        line = line.strip()

        if line.startswith("## Quotation"):
            # flush previous quote
            if table_data:
                subtotal = sum(p * q for p, q in price_qty_list)
                subtotal_str = f"{subtotal:,.0f}"
                build_table(current_quotation, table_data, bg_color=colors.HexColor("#4472C4"))
                elements.append(Paragraph(f"<b>Total Price ({current_quotation}):</b> ¥{subtotal_str}", styles["Normal"]))
                total_prices.append((current_quotation, subtotal_str))
                elements.append(Spacer(1, 10))

            current_quotation = line.replace("##", "").strip()
            table_data = [["Product Name", "Specs", "Price ($)", "Qty"]]
            price_qty_list = []
            inside_quote = True

        elif line.startswith("Product Name:"):
            pname = line.split(":", 1)[1].strip()

        elif line.startswith("Specs:"):
            specs = line.split(":", 1)[1].strip()

        elif line.startswith("Price:"):
            raw = line.split(":", 1)[1].strip()
            # strip out anything but digits and dot
            clean = re.sub(r"[^\d\.]", "", raw)
            price = float(clean)


        elif line.startswith("Quantity:"):
            qty = int(line.split(":", 1)[1].strip())
            table_data.append([pname, specs, f"{price:,.0f}", str(qty)])
            price_qty_list.append((price, qty))

        elif line.startswith("## Recommendation"):
            # flush last quote before recommendation
            if table_data and inside_quote:
                subtotal = sum(p * q for p, q in price_qty_list)
                subtotal_str = f"{subtotal:,.0f}"
                build_table(current_quotation, table_data, bg_color=colors.HexColor("#4472C4"))
                elements.append(Paragraph(f"<b>Total Price ({current_quotation}):</b> ¥{subtotal_str}", styles["Normal"]))
                total_prices.append((current_quotation, subtotal_str))
                elements.append(Spacer(1, 10))
            inside_quote = False
            elements.append(Paragraph("<b>🎯 Recommendation</b>", styles["Heading3"]))

        elif not inside_quote and line:
            recommendation_lines.append(line)

    # Summary tables if any left unflushed
    if table_data and inside_quote:
        build_table(current_quotation, table_data, bg_color=colors.HexColor("#4472C4"))

    # Pricing Summary
    elements.append(Spacer(1, 12))
    elements.append(Paragraph("<b>📊 Pricing Summary</b>", styles["Heading3"]))
    for qname, tprice in total_prices:
        elements.append(Paragraph(f"• {qname}: ¥{tprice}", styles["Normal"]))

    # Recommendation Section
    if recommendation_lines:
        elements.append(Spacer(1, 12))
        elements.append(Paragraph("<b>✅ Best Recommendation</b>", styles["Heading3"]))
        for line in recommendation_lines:
            elements.append(Paragraph(line, highlight_style))

    doc.build(elements)
    buffer.seek(0)

    with open("static/hardware_quotation.pdf", "wb") as f:
        f.write(buffer.getvalue())
    return buffer

# --- LangGraph App Logic ---
class State(TypedDict):
    question: str
    context: List[Document]
    answer: str

def retrieve(state: State):
    return {"context": vector_store.similarity_search(state["question"])}

def generate(state: State):
    docs_content = "\n\n".join(doc.page_content for doc in state["context"])
    messages = prompt.invoke({"question": state["question"], "context": docs_content})
    response = llm.invoke(messages)
    return {"answer": response.content}

graph_builder = StateGraph(State).add_sequence([retrieve, generate])
graph_builder.add_edge(START, "retrieve")
graph = graph_builder.compile()

# --- Streamlit UI ---
st.title("💻 Computer Hardware Sales Assistant")

if "result" not in st.session_state:
    st.session_state.result = None

st.header("--- Client Information ---")

client_name = st.text_input("Client Name", placeholder="e.g., John Doe")
client_company = st.text_input("Company Name", placeholder="e.g., ABC Corp")
client_email = st.text_input("Email", placeholder="e.g., john@example.com")
client_phone = st.text_input("Phone Number", placeholder="e.g., +1-234-567-890")

# Optionally store in session
st.session_state.client_info = {
    "name": client_name,
    "company": client_company,
    "email": client_email,
    "phone": client_phone
}

user_query = st.text_input("Enter your query:", placeholder="E.g., Best PC setup for video editing...")

if user_query and st.button("💬 Get Recommendation"):
    with st.spinner("Processing your query..."):
        st.session_state.result = graph.invoke({"question": user_query})

if st.session_state.result:
    st.subheader("💡 Suggested Answer")
    st.write(st.session_state.result["answer"])

    enc = tiktoken.encoding_for_model("gpt-4")
    input_tokens = len(enc.encode(user_query))
    output_tokens = len(enc.encode(st.session_state.result["answer"]))
    st.markdown(f"🔢 Input tokens: {input_tokens} | Output tokens: {output_tokens} | Total: {input_tokens + output_tokens}")

    pdf_bytes = generate_pdf(st.session_state.result["answer"], st.session_state.client_info)
    st.download_button(
        label="📄 Download Quotation as PDF",
        data=pdf_bytes,
        file_name="hardware_quotation.pdf",
        mime="application/pdf"
    )

    slide_bytes = generate_slides(st.session_state.result["answer"], st.session_state.client_info)
    st.download_button(
        label="📊 Download Slides (PPTX)",
        data=slide_bytes,
        file_name="hardware_quotation.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

    st.markdown("### 📤 Send Quotation via Email")

    send_email_btn = st.button("Send Email")
    if send_email_btn:
        if not client_email:
            st.error("Please enter the client's email above.")
        else:
            with st.spinner("Sending email..."):
                success, msg = send_email_with_attachment(
                    to_email=client_email,
                    subject="Hardware Quotation from Otsuka Corporation",
                    body=f"Dear {client_name},\n\nPlease find attached your hardware quotation.\n\nRegards,\nOtsuka Corporation"
                    #attachment_bytes=pdf_bytes,
                    #filename="hardware_quotation.pdf"
                )
                st.success(msg) if success else st.error(msg)

    pdf_path = "static/hardware_quotation.pdf"
    with open(pdf_path, "wb") as f:
        f.write(pdf_bytes.getvalue())

# Construct URL (adjust this to your deployed Streamlit domain)
    base_url = os.getenv("APP_BASE_URL", "http://localhost:8501")  # Override in prod
    pdf_url = f"{base_url}/static/hardware_quotation.pdf"

