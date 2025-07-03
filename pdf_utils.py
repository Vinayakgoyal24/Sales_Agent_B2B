# import os
# import time
# import re
# import pandas as pd
# import tiktoken
# from dotenv import load_dotenv
# from io import BytesIO
# import streamlit as st
# from langchain_community.document_loaders import TextLoader
# from langchain_core.documents import Document
# from langchain_text_splitters import RecursiveCharacterTextSplitter
# from langchain_chroma import Chroma
# from langchain_openai import AzureOpenAIEmbeddings, AzureChatOpenAI
# from langchain.prompts import ChatPromptTemplate
# from langgraph.graph import START, StateGraph
# from typing_extensions import List, TypedDict
# from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image
# from reportlab.lib.pagesizes import A4
# from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
# from reportlab.lib import colors
# from reportlab.lib.enums import TA_CENTER
# from reportlab.pdfbase.pdfmetrics import stringWidth
# from datetime import datetime, timedelta
# from pptx import Presentation
# from pptx.util import Inches, Pt
# from pptx.enum.shapes import MSO_SHAPE
# from pptx.dml.color import RGBColor
# from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
# from PIL import Image as PILImage
# import io
# import smtplib
# from email.mime.text import MIMEText
# from email.mime.multipart import MIMEMultipart
# from email.mime.base import MIMEBase
# from email import encoders
# import pandas as pd
# import base64
# import tiktoken
# from pptx import Presentation
# from pptx.util import Inches, Pt
# from pptx.dml.color import RGBColor
# from pptx.enum.text import PP_ALIGN
# from io import BytesIO
# import os
# import re

# # --- STT and TTS libraries ---
# from io import BytesIO
# import threading
# import tempfile
# import torch
# from faster_whisper import WhisperModel
# from streamlit_mic_recorder import mic_recorder
# from gtts import gTTS
# from reportlab.pdfbase import pdfmetrics
# from reportlab.pdfbase.ttfonts import TTFont


# def generate_pdf(quotation_text: str, client_info: dict) -> BytesIO:
#     startpdf=time.time()

#     # ── 1️⃣ language flag & labels ─────────────────────────────────────────
#     JA = bool(re.search(r"[ぁ-んァ-ヶ一-龯]", quotation_text))

#     TXT = {
#         "company":   "大塚商会"                if JA else "Otsuka Shokai",
#         "addr":      "〒102-8573 東京都千代田区飯田橋 2-18-4"
#                       if JA else "Head Office, 2-18-4 Iidabashi, Chiyoda-ku, Tokyo 102-8573",
#         "web":       "www.otsuka-shokai.co.jp",
#         "quotation": "–– 見積もり ––"          if JA else "--- Quotation ---",
#         "client":    "–– 顧客情報 ––"          if JA else "--- Client Information ---",
#         "date":      "発行日"                  if JA else "Date of Issue",
#         "valid":     "有効期限 (7日間)"        if JA else "Validity (7 days)",
#         "labels": {
#             "name":  "クライアント名"          if JA else "Client Name",
#             "company":"会社名"                if JA else "Client Company",
#             "email": "メール"                 if JA else "Email",
#             "phone": "電話番号"               if JA else "Phone",
#         },
#         "headers":   ["製品名","スペック","価格","数量"] if JA
#                      else ["Product Name","Specs","Price ($)","Qty"],
#         "total":     "合計価格"               if JA else "Total Price",
#         "recosec":   "推薦"                 if JA else "Recommendation",
#         "summary":   "価格サマリー"         if JA else "Pricing Summary",
#         "best":      "ベスト推薦"           if JA else "Best Recommendation",
#     }

#     # ── 2️⃣ register JP font if needed ────────────────────────────────────
#     base_font = "Helvetica"
#     if JA:
#         font_path = os.path.join("fonts", "NotoSansCJKjp-Regular.otf")  # ensure file exists
#         if os.path.exists(font_path):
#             pdfmetrics.registerFont(TTFont("NotoSansJP", font_path))
#             base_font = "NotoSansJP"
#         else:
#             print("[WARN] Japanese font not found; JP text may appear as □□")

#     buffer = BytesIO()
#     doc = SimpleDocTemplate(buffer,
#                             pagesize=A4,
#                             rightMargin=30, leftMargin=30,
#                             topMargin=30, bottomMargin=30)
#     styles = getSampleStyleSheet()
#     title_style = ParagraphStyle(name="TitleCenter", parent=styles["Title"],
#                                  alignment=TA_CENTER, fontName=base_font)
#     highlight_style = ParagraphStyle(name="Highlight", parent=styles["Normal"],
#                                      fontName=base_font, textColor=colors.black, fontSize=12)
#     normal_font = ParagraphStyle(name="NormalJP", parent=styles["Normal"],
#                                  fontName=base_font)
#     elements = []

#     # --- Header: Company Info ---
#     logo_path = "otsuka_im.png"
#     if os.path.exists(logo_path):
#         logo = Image(logo_path, width=100, height=50)
#         logo.hAlign = "LEFT"
#         elements.append(logo)

#     elements.append(Paragraph(TXT["company"], title_style))
#     elements.append(Paragraph(TXT["addr"], normal_font))
#     elements.append(Paragraph(f"Website: <a href='https://{TXT['web']}/'>{TXT['web']}</a>",
#                               normal_font))

#     elements.append(Spacer(1, 12))

#     # --- Quotation Info ---
#     elements.append(Paragraph(f"<b>{TXT['quotation']}</b>", styles["Heading2"]))
#     today = datetime.now()
#     validity = today + timedelta(days=7)
#     elements.append(Paragraph(f"{TXT['date']}: {today.strftime('%Y-%m-%d')}", normal_font))
#     elements.append(Paragraph(f"{TXT['valid']}: {validity.strftime('%Y-%m-%d')}", normal_font))

#     elements.append(Spacer(1, 12))

#     # --- Client Info (Dummy) ---
#     elements.append(Paragraph(f"<b>{TXT['client']}</b>", styles["Heading3"]))
#     client_info_lines = [
#         f"1. {TXT['labels']['name']}: {client_info.get('name','')}",
#         f"2. {TXT['labels']['company']}: {client_info.get('company','')}",
#         f"3. {TXT['labels']['email']}: {client_info.get('email','')}",
#         f"4. {TXT['labels']['phone']}: {client_info.get('contact','')}",
#     ]
#     for line in client_info_lines:
#         elements.append(Paragraph(line, normal_font))

#     elements.append(Spacer(1, 12))

#     # --- Parse LLM output and build tables ---
#     lines = quotation_text.strip().splitlines()
#     table_data = []
#     total_prices = []
#     price_qty_list = []
#     recommendation_lines = []
#     current_quotation = ""
#     inside_quote = False

#     def build_table(title, data, bg_color):
#         table_style = styles = getSampleStyleSheet()
#         cell_style = styles["Normal"]
#         font_name = base_font

#         font_size = 10

#         # Convert data cells to Paragraphs (except header)
#         tbl = [data[0]]
#         for row in data[1:]:
#             tbl.append([Paragraph(str(cell), cell_style) for cell in row])

#         # Dynamic column widths
#         transposed = list(zip(*data))
#         col_widths = []
#         for col in transposed:
#             max_w = max(stringWidth(str(item), font_name, font_size) for item in col)
#             col_widths.append(min(max_w + 20, 200))

#         table = Table(tbl, hAlign="LEFT", colWidths=col_widths)
#         table.setStyle(TableStyle([
#             ("BACKGROUND", (0, 0), (-1, 0), bg_color),
#             ("TEXTCOLOR", (0, 0), (-1, 0), colors.whitesmoke),
#             ("GRID", (0, 0), (-1, -1), 0.7, colors.black),
#             ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
#             ("FONTSIZE", (0, 0), (-1, -1), font_size),
#             ("ALIGN", (2, 1), (2, -1), "RIGHT"),
#             ("ALIGN", (3, 1), (3, -1), "CENTER"),
#             ("VALIGN", (0, 1), (-1, -1), "TOP"),
#         ]))

#         elements.append(Paragraph(f"<b>{title}</b>", styles["Heading4"]))
#         elements.append(table)

#     for line in lines:
#         line = line.strip()

#         if line.startswith("## Quotation"):
#             # flush previous quote
#             if table_data:
#                 subtotal = sum(p * q for p, q in price_qty_list)
#                 subtotal_str = f"{subtotal:,.0f}"
#                 build_table(current_quotation, table_data, bg_color=colors.HexColor("#4472C4"))
#                 elements.append(Paragraph(f"<b>{TXT['total']} ({current_quotation}):</b> ¥{subtotal_str}", normal_font))
#                 total_prices.append((current_quotation, subtotal_str))
#                 elements.append(Spacer(1, 10))

#             current_quotation = line.replace("##", "").strip()
#             table_data = [["Product Name", "Specs", "Price ($)", "Qty"]]
#             price_qty_list = []
#             inside_quote = True

#         if line.startswith("Product Name:"):
#             pname = line.split(":", 1)[1].strip()
#             specs, price, qty = "", 0.0, 0  # reset

#         elif line.startswith("Specs:"):
#             specs = line.split(":", 1)[1].strip()

#         elif line.startswith("Price:"):
#             raw = line.split(":", 1)[1].strip()
#             clean = re.sub(r"[^\d\.]", "", raw)
#             price = float(clean) if clean else 0.0

#         elif line.startswith("Quantity:"):
#             qty_raw = line.split(":", 1)[1].strip()
#             qty = int(qty_raw) if qty_raw.isdigit() else 0
#     # append only when all fields have been collected
#             table_data.append([pname, specs, f"{price:,.0f}", str(qty)])
#             price_qty_list.append((price, qty))

#         elif line.startswith("## Recommendation"):
#             # flush last quote before recommendation
#             if table_data and inside_quote:
#                 subtotal = sum(p * q for p, q in price_qty_list)
#                 subtotal_str = f"{subtotal:,.0f}"
#                 build_table(current_quotation, table_data, bg_color=colors.HexColor("#4472C4"))
#                 elements.append(Paragraph(f"<b>{TXT['total']} ({current_quotation}):</b> ¥{subtotal_str}", normal_font))
#                 total_prices.append((current_quotation, subtotal_str))
#                 elements.append(Spacer(1, 10))
#             inside_quote = False
#             elements.append(Paragraph(f"<b>{TXT['recosec']}</b>", styles["Heading3"]))

#         elif not inside_quote and line:
#             recommendation_lines.append(line)

#     # Summary tables if any left unflushed
#     if table_data and inside_quote:
#         build_table(current_quotation, table_data, bg_color=colors.HexColor("#4472C4"))

#     # Pricing Summary
#     elements.append(Spacer(1, 12))
#     elements.append(Paragraph(f"<b>{TXT['summary']}</b>", styles["Heading3"]))
#     for qname, tprice in total_prices:
#         elements.append(Paragraph(f"• {qname}: ¥{tprice}", styles["Normal"]))

#     # Recommendation Section
#     if recommendation_lines:
#         elements.append(Spacer(1, 12))
#         elements.append(Paragraph(f"<b>{TXT['best']}</b>", styles["Heading3"]))
#         for line in recommendation_lines:
#             elements.append(Paragraph(line, highlight_style))

#     doc.build(elements)
#     buffer.seek(0)

#     with open("static/hardware_quotation.pdf", "wb") as f:
#         f.write(buffer.getvalue())
#     elapsedp=time.time()-startpdf
#     print(f"pdf download: {elapsedp}")
#     return buffer



import os
import time
import re
import pandas as pd
import tiktoken
from dotenv import load_dotenv
from io import BytesIO
import streamlit as st
from langchain_community.document_loaders import TextLoader
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_chroma import Chroma
from langchain_openai import AzureOpenAIEmbeddings, AzureChatOpenAI
from langchain.prompts import ChatPromptTemplate
from langgraph.graph import START, StateGraph
from typing_extensions import List, TypedDict
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER
from reportlab.pdfbase.pdfmetrics import stringWidth
from datetime import datetime, timedelta
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from PIL import Image as PILImage
import io
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from email.mime.base import MIMEBase
from email import encoders
import pandas as pd
import base64
import tiktoken
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from io import BytesIO
import os
import re
from reportlab.platypus import PageBreak
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.lib.enums import TA_LEFT

pdfmetrics.registerFont(TTFont("NotoSansJP-Regular", "statics/fonts/NotoSansJP-VariableFont_wght.ttf"))
pdfmetrics.registerFont(TTFont("NotoSansJP", "statics/fonts/NotoSansJP-Bold.ttf"))

def contains_japanese(text):
    return bool(re.search(r"[\u3040-\u30FF\u4E00-\u9FFF]", text))

def generate_pdf(quotation_text: str, client_info: dict) -> BytesIO:
    startpdf=time.time()
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer,
                            pagesize=A4,
                            rightMargin=30, leftMargin=30,
                            topMargin=30, bottomMargin=30)
    styles = getSampleStyleSheet()

    has_japanese= contains_japanese(quotation_text)
    font_regular = "NotoSansJP-Regular" if has_japanese else "Helvetica"
    font_bold = "NotoSansJP" if has_japanese else "Helvetica-Bold"
    # Styles
    title_style = ParagraphStyle(name="TitleCenter", parent=styles["Title"], alignment=TA_CENTER, fontName=font_bold)
    normal_style = ParagraphStyle(name="NormalText", parent=styles["Normal"], fontName=font_regular, fontSize=12, leading=12, alignment=TA_LEFT)
    highlight_style = ParagraphStyle(name="Highlight", parent=styles["Normal"], fontName=font_regular, textColor=colors.black, fontSize=12)
    heading2_style = ParagraphStyle(name="Heading2", parent=styles["Heading2"], fontName=font_bold)
    heading3_style = ParagraphStyle(name="Heading3", parent=styles["Heading3"], fontName=font_bold)
    heading4_style = ParagraphStyle(name="Heading4", parent=styles["Heading4"], fontName=font_bold)


    elements = []

    # --- Header: Company Info ---
    logo_path = "otsuka_im.png"
    if os.path.exists(logo_path):
        logo = Image(logo_path, width=100, height=50)
        logo.hAlign = "LEFT"
        elements.append(logo)

    elements.append(Paragraph("Otsuka Shokai", title_style))
    elements.append(Paragraph("Head Office, 2-18-4 Iidabashi, Chiyoda-ku, Tokyo 102-8573", styles["Normal"]))
    elements.append(Paragraph("Website: <a href='https://www.otsuka-shokai.co.jp/'>www.otsuka-shokai.co.jp</a>",
                              styles["Normal"]))
    elements.append(Spacer(1, 12))

    # --- Quotation Info ---
    header_text = "--- 見積もり ---" if has_japanese else "--- Quotation ---"
    elements.append(Paragraph(f"<b>{header_text}</b>", heading2_style))
    today = datetime.now()
    validity = today + timedelta(days=7)
    elements.append(Paragraph(f"{'1. 発行日' if has_japanese else '1. Date of Issue'}: {today.strftime('%Y-%m-%d')}", normal_style))
    elements.append(Paragraph(f"{'2. 有効期限' if has_japanese else '2. Validity'}: {validity.strftime('%Y-%m-%d')} (7 days)", normal_style))

    # --- Client Info (Dummy) ---
    client_header = "--- クライアント情報 ---" if has_japanese else "--- Client Information ---"
    elements.append(Paragraph(f"<b>{client_header}</b>", heading3_style))
    if has_japanese:
        client_info_lines = [
        f"<b>1. クライアント名:</b> {client_info.get('name', '')}",
        f"<b>2. 会社名:</b> {client_info.get('company', '')}",
        f"<b>3. メール:</b> {client_info.get('email', '')}",
        f"<b>4. 電話番号:</b> {client_info.get('contact', '')}",
    ]
    else:
        client_info_lines = [
        f"<b>1. Client Name:</b> {client_info.get('name', '')}",
        f"<b>2. Client Company:</b> {client_info.get('company', '')}",
        f"<b>3. Email:</b> {client_info.get('email', '')}",
        f"<b>4. Phone:</b> {client_info.get('contact', '')}",
    ]

    for line in client_info_lines:
        elements.append(Paragraph(line, normal_style))
    elements.append(Spacer(1, 12))

    # --- Parse LLM output and build tables ---
    lines = quotation_text.strip().splitlines()
    table_data = []
    total_prices = []
    price_qty_list = []
    recommendation_lines = []
    current_quotation = ""
    inside_quote = False

    def build_table(title, data, bg_color):
        tbl = [data[0]] + [[Paragraph(str(cell), normal_style) for cell in row] for row in data[1:]]
        col_widths = [min(max(stringWidth(str(item), font_regular, 10) + 20 for item in col), 200) for col in zip(*data)]
        table = Table(tbl, hAlign="LEFT", colWidths=col_widths)
        table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), bg_color),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.whitesmoke),
            ("GRID", (0, 0), (-1, -1), 0.7, colors.black),
            ("FONTNAME", (0, 0), (-1, -1), font_regular),
            ("FONTSIZE", (0, 0), (-1, -1), 10),
            ("ALIGN", (2, 1), (2, -1), "RIGHT"),
            ("ALIGN", (3, 1), (3, -1), "CENTER"),
            ("VALIGN", (0, 1), (-1, -1), "TOP"),
        ]))
        elements.append(Paragraph(f"<b>{title}</b>", heading4_style))
        elements.append(table)

    for line in lines:
        line = line.strip()
        if line.startswith("## Quotation") or line.startswith("## 見積もり"):
            if table_data:
                subtotal = sum(p * q for p, q in price_qty_list)
                build_table(current_quotation, table_data, colors.HexColor("#4472C4"))
                elements.append(Paragraph(f"<b>- Total Price ({current_quotation}):</b> ¥{subtotal:,.0f}", normal_style))
                elements.append(Spacer(1, 10))
                total_prices.append((current_quotation, f"{subtotal:,.0f}"))
            current_quotation = line.replace("##", "").strip()
            table_data = [["商品名", "仕様", "価格", "数量"]] if has_japanese else [["Product Name", "Specs", "Price", "Qty"]]
            price_qty_list = []
            inside_quote = True
        elif line.startswith("Product Name:") or line.startswith("商品名:"):
            pname = line.split(":", 1)[1].strip()
        elif line.startswith("Specs:") or line.startswith("仕様:"):
            specs = line.split(":", 1)[1].strip()
        elif line.startswith("Price:") or line.startswith("価格:"):
            price = float(re.sub(r"[^\d.]", "", line.split(":", 1)[1].strip()) or 0)
        elif line.startswith("Quantity:") or line.startswith("数量:"):
            qty = int(line.split(":", 1)[1].strip() or 0)
            table_data.append([pname, specs, f"{price:,.0f}", str(qty)])
            price_qty_list.append((price, qty))
        elif line.startswith("## Recommendation") or line.startswith("## 推奨案"):
            if table_data and inside_quote:
                subtotal = sum(p * q for p, q in price_qty_list)
                build_table(current_quotation, table_data, colors.HexColor("#4472C4"))
                elements.append(Paragraph(f"<b>Total Price ({current_quotation}):</b> ¥{subtotal:,.0f}", normal_style))
                elements.append(Spacer(1, 10))
                total_prices.append((current_quotation, f"{subtotal:,.0f}"))
            inside_quote = False
        elif not inside_quote and line:
            recommendation_lines.append(line)

    if table_data and inside_quote:
        subtotal = sum(p * q for p, q in price_qty_list)
        build_table(current_quotation, table_data, colors.HexColor("#4472C4"))
        elements.append(Paragraph(f"<b>Total Price ({current_quotation}):</b> ¥{subtotal:,.0f}", normal_style))
        total_prices.append((current_quotation, f"{subtotal:,.0f}"))

    # Summary & Recommendation
    elements.append(Spacer(1, 12))
    elements.append(Paragraph(f"<b>{'--- 価格まとめ ---' if has_japanese else '--- Pricing Summary ---'}</b>", heading3_style))
    for qname, tprice in total_prices:
        elements.append(Paragraph(f"• {qname}: ¥{tprice}", normal_style))

    if recommendation_lines:
        elements.append(PageBreak())
        elements.append(Spacer(1, 12))
        rec_header = "--- ベスト推奨案 ---" if contains_japanese(quotation_text) else "--- Best Recommendation ---"
        elements.append(Paragraph(f"<b>{rec_header}</b>", title_style))

        for line in recommendation_lines:
            if contains_japanese(line):
                elements.append(Paragraph(line, highlight_style))
            else:
                elements.append(Paragraph(line, normal_style))


    doc.build(elements)
    buffer.seek(0)
    with open("static/hardware_quotation.pdf", "wb") as f:
        f.write(buffer.getvalue())
    print(f"pdf download: {time.time() - startpdf:.2f} seconds")
    return buffer