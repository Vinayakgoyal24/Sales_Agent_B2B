import os
import re
import pandas as pd
import tiktoken
from dotenv import load_dotenv
from io import BytesIO
import streamlit as st
from langchain_community.document_loaders import TextLoader
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_chroma import Chroma
from langchain_openai import AzureOpenAIEmbeddings, AzureChatOpenAI
from langchain.prompts import ChatPromptTemplate
from langgraph.graph import START, StateGraph
from typing_extensions import List, TypedDict
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER
from reportlab.pdfbase.pdfmetrics import stringWidth
from datetime import datetime, timedelta
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from PIL import Image as PILImage
import io
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from email.mime.base import MIMEBase
from email import encoders
import pandas as pd
import base64
import tiktoken
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from io import BytesIO
import os
import re

# --- STT and TTS libraries ---
from io import BytesIO
import threading
import tempfile
import torch
from faster_whisper import WhisperModel
from streamlit_mic_recorder import mic_recorder
from gtts import gTTS


# Load env vars
load_dotenv()

# --- Initialize Azure Clients ---
llm = AzureChatOpenAI(
    azure_endpoint=os.getenv("AZURE_OPENAI_ENDPOINT"),
    azure_deployment=os.getenv("AZURE_OPENAI_DEPLOYMENT_NAME"),
    openai_api_version=os.getenv("AZURE_OPENAI_API_VERSION"),
)

embeddings = AzureOpenAIEmbeddings(
    azure_endpoint=os.getenv("AZURE_OPENAI_EMBEDDING_ENDPOINT"),
    azure_deployment=os.getenv("AZURE_OPENAI_EMBEDDING_DEPLOYMENT_NAME"),
    openai_api_version=os.getenv("AZURE_OPENAI_EMBEDDING_API_VERSION"),
)

vector_store = Chroma(
    collection_name="example_collection",
    embedding_function=embeddings,
    persist_directory="./chroma_langchain_db",
)


LANGUAGES = {"English": "en", "日本語": "ja"}
DEFAULT_LANG = "English"

TRANSLATE = {
    "title":            {"en": "💻 Computer Hardware Sales Assistant",
                         "ja": "💻 コンピュータハードウェア販売アシスタント"},
    "help_header":      {"en": "How can I help you today? 😊",
                         "ja": "本日はいかがいたしましょうか？ 😊"},
    "query_placeholder": {"en": "E.g., Best PC setup for video editing…",
                          "ja": "例: 動画編集に最適な PC セットアップ…"},
    "get_reco":         {"en": "💬 Get Recommendation",
                         "ja": "💬 見積もりを取得"},
    "sample1":          {"en": "💼 Lightweight business laptop",
                         "ja": "💼 軽量ビジネスノート"},
    "sample2":          {"en": "🎮 Good gaming mouse",
                         "ja": "🎮 高性能ゲーミングマウス"},
    "sample3":          {"en": "📸 4K webcam for reels",
                         "ja": "📸 リール用 4K ウェブカメラ"},
    "quotation_header": {"en": "💡Quotation tailored to your needs",
                         "ja": "💡 お客様のニーズに合わせたお見積もり"},
    "client_expander":  {"en": "📋 Fill Client Information",
                         "ja": "📋 顧客情報を入力"},
    "name_ph":          {"en": "e.g., John Doe",
                         "ja": "例: 田中 太郎"},
    "email_ph":         {"en": "e.g., john@example.com",
                         "ja": "例: taro@example.com"},
    "company_ph":       {"en": "e.g., ABC Corp",
                         "ja": "例: 株式会社ABC"},
    "phone_ph":         {"en": "e.g., +1-234-567-890",
                         "ja": "例: 03-1234-5678"},
    "speak_btn":        {"en": "🔊 Speak this quotation",
                         "ja": "🔊 見積もりを読み上げ"},
    "export_pdf":       {"en": "📄 Download PDF",
                         "ja": "📄 PDF をダウンロード"},
    "export_ppt":       {"en": "📊 Download PPTX",
                         "ja": "📊 PPTX をダウンロード"},
    "email_quote":      {"en": "📧 Email Quotation",
                         "ja": "📧 見積もりをメール送信"},
    "email_err":        {"en": "Please enter the client's email above.",
                         "ja": "上記に顧客のメールアドレスを入力してください。"},
    "email_body":       {"en": "Dear {name},\n\nPlease find attached your hardware quotation.\n\nRegards,\nOtsuka Shokai",
                         "ja": "{name} 様\n\nハードウェアのお見積書を添付いたしましたのでご確認ください。\n\nよろしくお願いいたします。\n大塚商会"},
    "feedback_header":  {"en": "Any changes or feedback?",
                         "ja": "修正点やご要望はございますか？"},
    "feedback_placeholder": {"en": "Suggest changes or type 'thanks' to finalize:",
                             "ja": "修正内容を入力、もしくは「thanks」と入力して確定してください:"},
    "save_chat":        {"en": "💾 Save This Chat Session",
                         "ja": "💾 このチャットを保存"},
    "no_chat_warning":  {"en": "No chat history to save.",
                         "ja": "保存するチャット履歴がありません。"},
    "session_saved":    {"en": "Session saved as {file}",
                         "ja": "{file} として保存しました"},
    "restart":          {"en": "🔄 Restart Session",
                         "ja": "🔄 セッションをリスタート"},
    "export_header":    {"en": "### 📎 Export Options",
                         "ja": "### 📎 エクスポートオプション"},
    "submit_feedback":  {"en": "🔁 Submit Feedback",
                         "ja": "🔁 フィードバック送信"},
    "revised_quotation":    {"en": "💡 Revised Quotation",
                             "ja": "💡 修正後の見積もり"},
    "transcribing":     {"en": "Transcribing…",
                         "ja": "文字起こし中…"},
    "transcribing_fb":  {"en": "Transcribing feedback…",
                         "ja": "フィードバックを文字起こし中…"},
    "finding":          {"en": "Finding the best products for you",
                         "ja": "最適な製品を探しています"},
    "audio_wait":       {"en": "Audio is still being prepared – please try again in a moment.",
                         "ja": "音声を準備中です。しばらくしてからお試しください。"},
    
    "sending_email":     {"en": "Sending email...",       "ja": "メール送信中…"},
    "total_lbl":         {"en": "Total: ${val}",          "ja": "合計: ¥{val}"},
    "best_reco":         {"en": "Best Recommendation",    "ja": "最適な提案"},
    "client_info_hdr":   {"en": "Client Information",     "ja": "顧客情報"},
    "quotation_hdr":     {"en": "Hardware Configuration Quotations",
                          "ja": "ハードウェア構成見積書"},
    "thank_you":         {"en": "Thank You",              "ja": "ありがとうございます"},

    "sidebar_sessions":      {"en": "📁 Chat Sessions",          "ja": "📁 チャットセッション"},
    "load_session":          {"en": "Load previous session",     "ja": "以前のセッションを読み込む"},
    "loaded_data_hdr":       {"en": "### 🗂 Loaded Session Data","ja": "### 🗂 読み込んだセッション"},
    # runtime & spinners
    "indexing_docs":         {"en": "Indexing documents...",     "ja": "ドキュメントをインデックス中…"},
    "incorporating":         {"en": "Incorporating your feedback...",
                              "ja": "フィードバックを反映中…"},
    "thank_finalized":       {"en": "🎉 Thank you! Quotation finalized.",
                              "ja": "🎉 ありがとうございます！ 見積もりを確定しました。"},
    # PDF/PPT headings left over
    "recommendation_hdr":    {"en": "Recommendation",            "ja": "おすすめ"},
    "pricing_summary":       {"en": "Pricing Summary",           "ja": "価格サマリー"},
    "total_price":           {"en": "Total Price",               "ja": "合計金額"},
    # footer / slide
    "footer_by":             {"en": "by Otsuka Shokai",          "ja": "大塚商会"},
    "page_lbl":              {"en": "Page {idx}",                "ja": "ページ {idx}"}


}

def t(key: str) -> str:
    """Return the UI string for the current language."""
    code = st.session_state.get("lang_code", LANGUAGES[DEFAULT_LANG])
    return TRANSLATE[key][code]

def _l(en: str, ja: str) -> str:
    """Return ja when UI is Japanese, else en – used inside PDF/PPT helpers."""
    return ja if st.session_state.get("lang_code", "en") == "ja" else en


# --- Load CSVs and Create Documents ---
def load_csv_as_documents(folder_path="data"):
    documents = []
    for file_name in os.listdir(folder_path):
        if file_name.endswith(".csv"):
            df = pd.read_csv(os.path.join(folder_path, file_name)).head(20)
            for _, row in df.iterrows():
                content = "\n".join([f"{k}: {v}" for k, v in row.items()])
                documents.append(Document(page_content=content, metadata={"source": file_name}))
    return documents

if vector_store._collection.count() == 0:
    with st.spinner(t("indexing_docs")):
        docs = load_csv_as_documents("data")
        splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        splits = splitter.split_documents(docs)
        vector_store.add_documents(splits)
        vector_store.persist()

# ────────────────────────────────────────────────────────────────────────────────
# Prompt Templates (EN / JA)
# ────────────────────────────────────────────────────────────────────────────────
SYSTEM_PROMPT_EN = (
    "You are a professional hardware sales assistant at Otsuka Shokai. "
    "Based on the user's request and the context, provide 2–3 detailed hardware configuration quotations. "
    "Each quotation should include:\n"
    "- Product Name\n- Specs\n- Price\n- Quantity\n- Total Price\n\n"
    "Use this structure:\n"
    "## Quotation 1\nProduct Name: ...\nSpecs: ...\nPrice: ...\nQuantity: ...\n...\nTotal Price: ...\n"
    "## Quotation 2 ...\n\n"
    "Then provide a clear comparison of the quotations and recommend the best one based on:\n"
    "- Price\n- Suitability for the user's need\n- Performance vs cost.\n"
    "Use a section titled:\n"
    "## Recommendation\n"
    "Mention why the chosen quote is the best and highlight key differences with others.\n\n"
    "Keep tone professional and brief. Do not fabricate information if context is insufficient."
)

SYSTEM_PROMPT_JA = (
    "あなたは大塚商会のプロフェッショナルなハードウェア営業アシスタントです。"
    "ユーザーの要望と文脈に基づき、2〜3件の詳細なハードウェア構成見積もりを提示してください。\n"
    "各見積もりには以下を含めます:\n"
    "- 製品名\n- 仕様\n- 価格\n- 数量\n- 合計金額\n\n"
    "フォーマット例:\n"
    "## Quotation 1\nProduct Name: ...\nSpecs: ...\nPrice: ...\nQuantity: ...\n...\nTotal Price: ...\n"
    "## Quotation 2 ...\n\n"
    "最後に見積もりを比較し、以下に基づき最適なものを推薦してください:\n"
    "- 価格\n- ユーザー要件への適合性\n- パフォーマンスとコストのバランス\n"
    "セクションタイトル:\n"
    "## Recommendation\n"
    "選定理由と他案との主な相違点を簡潔に述べてください。\n\n"
    "プロフェッショナルかつ簡潔な口調を維持し、文脈が不足している場合は情報を捏造しないでください。"
)

def build_prompt() -> ChatPromptTemplate:
    if st.session_state.lang_code == "ja":
        sys_msg  = SYSTEM_PROMPT_JA + "\n\n回答は必ず日本語で。"
        human_fmt = "質問: {question}\n\nコンテキスト:\n{context}"
    else:
        sys_msg  = SYSTEM_PROMPT_EN
        human_fmt = "Question: {question}\n\nContext:\n{context}"
    return ChatPromptTemplate.from_messages([
        ("system", sys_msg),
        ("human",  human_fmt)
    ])

def send_email_with_attachment(to_email: str, subject: str, body: str):

    smtp_port = 587
    smtp_server = "smtp.gmail.com"
    sender_email = "vinayak.otsuka@gmail.com"
    pswd = "djjvyfubleftjmwh"

    if not sender_email or not pswd:
        return False, "Missing sender email or password in environment variables."

    msg = MIMEMultipart()
    msg['From'] = sender_email
    msg['To'] = to_email
    msg['Subject'] = subject

    msg.attach(MIMEText(body, 'plain'))

    filename = "static/hardware_quotation.pdf"
    try:
        with open(filename, 'rb') as attachment:
            attachment_package = MIMEBase('application', 'octet-stream')
            attachment_package.set_payload(attachment.read())
            encoders.encode_base64(attachment_package)
            attachment_package.add_header('Content-Disposition', f"attachment; filename={filename}")
            msg.attach(attachment_package)
    except FileNotFoundError:
        return False, f"Attachment file not found: {filename}"

    text = msg.as_string()

    try:
        print("Connecting to server...")
        TIE_server = smtplib.SMTP(smtp_server, smtp_port)
        TIE_server.starttls()
        TIE_server.login(sender_email, pswd)
        print("Successfully connected to server")

        TIE_server.sendmail(sender_email, to_email, text)
        TIE_server.quit()
        return True, f"Email successfully sent to {to_email}"
    except Exception as e:
        return False, f"Failed to send email to {to_email}. Error: {e}"

def generate_slides(quotation_text: str, client_info: dict) -> BytesIO:
    prs = Presentation()
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    bg_color = RGBColor(245, 245, 245)  # Light Gray Background

    def set_slide_bg_color(slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    def add_logo(slide):
        if os.path.exists("otsuka_im.png"):
            left = slide_width - Inches(1.5)
            top = Inches(0.3)
            slide.shapes.add_picture("otsuka_im.png", left, top, width=Inches(1.2))

    def remove_title_placeholder(slide):
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                sp = shape
                sp.element.getparent().remove(sp.element)


    def add_footer(slide, index):
        margin_bottom = Inches(0.2)  # distance from bottom
        textbox_height = Inches(0.25)

    # "by Otsuka Shokai"
        footer = slide.shapes.add_textbox(
            Inches(0.3), slide_height - margin_bottom - textbox_height, Inches(3), textbox_height
        )
        p = footer.text_frame.add_paragraph()
        p.text = t("footer_by")
        p.font.size = Pt(10)
        p.font.italic = True
        p.font.color.rgb = RGBColor(120, 120, 120)

    # Page number
        page_num = slide.shapes.add_textbox(
            slide_width - Inches(1), slide_height - margin_bottom - textbox_height, Inches(0.7), textbox_height
        )
        p2 = page_num.text_frame.add_paragraph()
        p2.text = t("page_lbl").format(idx=index)
        p2.font.size = Pt(10)
        p2.font.color.rgb = RGBColor(120, 120, 120)


    def underline_title(slide):
        if slide.shapes.title:
            title_shape = slide.shapes.title
            p = title_shape.text_frame.paragraphs[0]
            p.font.underline = True
            p.font.color.rgb = RGBColor(0, 51, 102)
            p.font.size = Pt(28)

    def add_quotation_slide(prs, title, table_data, price_qty_list, index):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        set_slide_bg_color(slide)
        slide.shapes.title.text = title
        underline_title(slide)

        rows, cols = len(table_data), len(table_data[0])
        left, top, width = Inches(0.5), Inches(1.5), Inches(9)
        height = Inches(0.8 + 0.4 * rows)

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        for col_index in range(cols):
            table.columns[col_index].width = Inches(2.2)

        for r in range(rows):
            for c in range(cols):
                cell = table.cell(r, c)
                cell.text = str(table_data[r][c])
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(10)
                if r == 0:
                    para.font.bold = True
                    para.font.color.rgb = RGBColor(255, 255, 255)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0, 112, 192)
                else:
                    para.font.color.rgb = RGBColor(0, 0, 0)

        total = sum(p * q for p, q in price_qty_list)
        txBox = slide.shapes.add_textbox(Inches(6.5), top + height + Inches(0.3), Inches(2.5), Inches(0.8))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = t("total_lbl").format(val=f"{total:,.0f}")
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 112, 192)
        p.alignment = PP_ALIGN.RIGHT


        add_logo(slide)
        add_footer(slide, index)
        return slide

    slide_index = 1

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    remove_title_placeholder(slide)
    set_slide_bg_color(slide)
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    tf = box.text_frame
    p = tf.add_paragraph()
    p.text = t("quotation_hdr")   
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    subtitle = tf.add_paragraph()
    subtitle.text = _l(
    "Generated by Otsuka Shokai AI Sales Agent",
    "大塚商会 AI セールスエージェント生成")
    subtitle.font.size = Pt(18)
    subtitle.font.color.rgb = RGBColor(100, 100, 100)

    add_logo(slide)
    add_footer(slide, slide_index)
    slide_index += 1

    # Slide 2: Client Info
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    set_slide_bg_color(slide)
    slide.shapes.title.text = t("client_info_hdr")
    underline_title(slide)

    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8), Inches(4))
    tf = box.text_frame
    tf.word_wrap = True
    for key, en, ja in [
            ("name",    "Client Name",   "顧客名"),
            ("company", "Company",       "会社名"),
            ("email",   "Contact Email", "メールアドレス"),
            ("phone",   "Contact Phone", "電話番号")]:
        p = tf.add_paragraph()
        p.text = f"{_l(en, ja)}: {client_info.get(key,'')}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(50, 50, 50)

    add_logo(slide)
    add_footer(slide, slide_index)
    slide_index += 1

    # Parsing and Quotation Slides
    lines = quotation_text.strip().splitlines()
    current_quotation = ""
    table_data = []
    price_qty_list = []
    recommendation_lines = []
    inside_quote = False

    for line in lines:
        line = line.strip()
        if line.startswith("## Quotation"):
            if table_data:
                add_quotation_slide(prs, current_quotation, table_data, price_qty_list, slide_index)
                slide_index += 1
            current_quotation = line.replace("##", "").strip()
            table_data = [[ _l("Product Name","製品名"),
                _l("Specs","仕様"),
                _l("Price","価格"),
                _l("Qty","数量") ]]
            price_qty_list = []
            inside_quote = True
        elif line.startswith("Product Name:"):
            pname = line.split(":", 1)[1].strip()
        elif line.startswith("Specs:"):
            specs = line.split(":", 1)[1].strip()
        # elif line.startswith("Price:"):
        #     price = float(re.sub(r"[^\d.]", "", line.split(":", 1)[1].strip()))
        elif line.startswith("Price:"):
            raw = line.split(":", 1)[1].strip()
            clean = re.sub(r"[^\d\.]", "", raw)

            try:
                price = float(clean)
            except ValueError:
                price = 0.0

        elif line.startswith("Quantity:"):
            qty_raw = line.split(":", 1)[1].strip()
            digits = re.findall(r'\d+', qty_raw)
            if digits:
                qty = int(digits[0])
            else:
                print(f"[WARNING] Invalid quantity format '{qty_raw}' — defaulting to 1")
                qty = 1
            table_data.append([pname, specs, f"${price:,.0f}", str(qty)])
            price_qty_list.append((price, qty))
        elif line.startswith("## Recommendation"):
            if table_data:
                add_quotation_slide(prs, current_quotation, table_data, price_qty_list, slide_index)
                slide_index += 1
            table_data = []
            inside_quote = False
        elif not inside_quote and line:
            recommendation_lines.append(line)

    if table_data and inside_quote:
        add_quotation_slide(prs, current_quotation, table_data, price_qty_list, slide_index)
        slide_index += 1

    # Recommendation Slide
    if recommendation_lines:
        chunks = [recommendation_lines[i:i+10] for i in range(0, len(recommendation_lines), 10)]
        for idx, chunk in enumerate(chunks):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            set_slide_bg_color(slide)
            title = t("best_reco") + (f" (Part {idx+1})" if len(chunks) > 1 else "")
            slide.shapes.title.text = title
            underline_title(slide)

            box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(8.5), Inches(4.5))
            tf = box.text_frame
            tf.word_wrap = True
            font_size = Pt(20 if len(chunk) <= 6 else 16 if len(chunk) <= 10 else 14)
            for line in chunk:
                bullet = tf.add_paragraph()
                bullet.text = f"• {line}"
                bullet.level = 0
                bullet.font.size = font_size
                bullet.font.color.rgb = RGBColor(60, 60, 60)

            add_logo(slide)
            add_footer(slide, slide_index)
            slide_index += 1

    # Thank You Slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_bg_color(slide)
    remove_title_placeholder(slide)

    # Centered "Thank You"
    box = slide.shapes.add_textbox(Inches(2), Inches(1.8), Inches(6), Inches(1))
    p = box.text_frame.paragraphs[0]
    p.text = t("thank_you")
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.underline = True
    p.alignment = PP_ALIGN.CENTER

    # Centered contact info
    info_box = slide.shapes.add_textbox(Inches(1.5), Inches(3), Inches(7), Inches(2.5))
    tf = info_box.text_frame
    lines = [
        _l("We appreciate your interest in Otsuka Shokai.",
        "大塚商会にご関心をお寄せいただきありがとうございます。"),
        _l("For any inquiries, reach out at:",
        "ご不明点は下記までご連絡ください:"),
        "📧 support@otsuka-shokai.co.jp",
        "🌐 www.otsuka-shokai.co.jp"
    ]
    for line in lines:
        para = tf.add_paragraph()
        para.text = line
        para.font.size = Pt(18)
        para.alignment = PP_ALIGN.CENTER

    add_footer(slide, slide_index)


    pptx_buffer = BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    return pptx_buffer


def generate_pdf(quotation_text: str, client_info: dict) -> BytesIO:
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer,
                            pagesize=A4,
                            rightMargin=30, leftMargin=30,
                            topMargin=30, bottomMargin=30)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(name="TitleCenter", parent=styles["Title"], alignment=TA_CENTER)
    highlight_style = ParagraphStyle(name="Highlight", parent=styles["Normal"], textColor=colors.black, fontSize=12)

    elements = []

    # --- Header: Company Info ---
    logo_path = "otsuka_im.png"
    if os.path.exists(logo_path):
        logo = Image(logo_path, width=100, height=50)
        logo.hAlign = "LEFT"
        elements.append(logo)

    elements.append(Paragraph("Otsuka Shokai", title_style))
    elements.append(Paragraph("Head Office, 2-18-4 Iidabashi, Chiyoda-ku, Tokyo 102-8573", styles["Normal"]))
    elements.append(Paragraph("Website: <a href='https://www.otsuka-shokai.co.jp/'>www.otsuka-shokai.co.jp</a>",
                              styles["Normal"]))
    elements.append(Spacer(1, 12))

    # --- Quotation Info ---
    elements.append(Paragraph(f"<b>{_l('--- Quotation ---','--- 見積書 ---')}</b>", styles["Heading2"]))
    today = datetime.now()
    validity = today + timedelta(days=7)
    elements.append(Paragraph(f"Date of Issue: {today.strftime('%Y-%m-%d')}", styles["Normal"]))
    elements.append(Paragraph(f"Validity: {validity.strftime('%Y-%m-%d')} (7 days)", styles["Normal"]))
    elements.append(Spacer(1, 12))

    # --- Client Info (Dummy) ---
    elements.append(Paragraph(f"<b>{_l('--- Client Information ---','--- 顧客情報 ---')}</b>", styles["Heading3"]))
    client_info = [
        "1. Client Name: "+ client_info.get("name",""),
        "2. Client Company:"+ client_info.get("company",""),
        "3. Email: "+ client_info.get("email",""),
        "4. Phone: "+ client_info.get("phone","")
    ]
    for line in client_info:
        elements.append(Paragraph(line, styles["Normal"]))
    elements.append(Spacer(1, 12))

    # --- Parse LLM output and build tables ---
    lines = quotation_text.strip().splitlines()
    table_data = []
    total_prices = []
    price_qty_list = []
    recommendation_lines = []
    current_quotation = ""
    inside_quote = False

    def build_table(title, data, bg_color):
        table_style = styles = getSampleStyleSheet()
        cell_style = styles["Normal"]
        font_name = "Helvetica"
        font_size = 10

        # Convert data cells to Paragraphs (except header)
        tbl = [data[0]]
        for row in data[1:]:
            tbl.append([Paragraph(str(cell), cell_style) for cell in row])

        # Dynamic column widths
        transposed = list(zip(*data))
        col_widths = []
        for col in transposed:
            max_w = max(stringWidth(str(item), font_name, font_size) for item in col)
            col_widths.append(min(max_w + 20, 200))

        table = Table(tbl, hAlign="LEFT", colWidths=col_widths)
        table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), bg_color),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.whitesmoke),
            ("GRID", (0, 0), (-1, -1), 0.7, colors.black),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), font_size),
            ("ALIGN", (2, 1), (2, -1), "RIGHT"),
            ("ALIGN", (3, 1), (3, -1), "CENTER"),
            ("VALIGN", (0, 1), (-1, -1), "TOP"),
        ]))

        elements.append(Paragraph(f"<b>{title}</b>", styles["Heading4"]))
        elements.append(table)

    for line in lines:
        line = line.strip()

        if line.startswith("## Quotation"):
            # flush previous quote
            if table_data:
                subtotal = sum(p * q for p, q in price_qty_list)
                subtotal_str = f"{subtotal:,.0f}"
                build_table(current_quotation, table_data, bg_color=colors.HexColor("#4472C4"))
                elements.append(Paragraph(f"<b>Total Price ({current_quotation}):</b> ¥{subtotal_str}", styles["Normal"]))
                total_prices.append((current_quotation, subtotal_str))
                elements.append(Spacer(1, 10))

            current_quotation = line.replace("##", "").strip()
            table_data = [[ _l("Product Name","製品名"),
                _l("Specs","仕様"),
                _l("Price ($)","価格 ($)"),
                _l("Qty","数量") ]]
            price_qty_list = []
            inside_quote = True

        elif line.startswith("Product Name:"):
            pname = line.split(":", 1)[1].strip()

        elif line.startswith("Specs:"):
            specs = line.split(":", 1)[1].strip()

        elif line.startswith("Price:"):
            raw = line.split(":", 1)[1].strip()
            # strip out anything but digits and dot
            clean = re.sub(r"[^\d\.]", "", raw)
            price = float(clean)


        elif line.startswith("Quantity:"):
            qty = int(line.split(":", 1)[1].strip())
            table_data.append([pname, specs, f"{price:,.0f}", str(qty)])
            price_qty_list.append((price, qty))

        elif line.startswith("## Recommendation"):
            # flush last quote before recommendation
            if table_data and inside_quote:
                subtotal = sum(p * q for p, q in price_qty_list)
                subtotal_str = f"{subtotal:,.0f}"
                build_table(current_quotation, table_data, bg_color=colors.HexColor("#4472C4"))
                elements.append(Paragraph(f"<b>{t('total_price')} ({current_quotation}):</b> ¥{subtotal_str}", styles["Normal"]))
                total_prices.append((current_quotation, subtotal_str))
                elements.append(Spacer(1, 10))
            inside_quote = False
            elements.append(Paragraph("<b>🎯 "+t("recommendation_hdr")+"</b>", styles["Heading3"]))

        elif not inside_quote and line:
            recommendation_lines.append(line)

    # Summary tables if any left unflushed
    if table_data and inside_quote:
        build_table(current_quotation, table_data, bg_color=colors.HexColor("#4472C4"))

    # Pricing Summary
    elements.append(Spacer(1, 12))
    elements.append(Paragraph("<b>📊 "+t('pricing_summary')+"</b>", styles["Heading3"]))
    for qname, tprice in total_prices:
        elements.append(Paragraph(f"• {qname}: ¥{tprice}", styles["Normal"]))

    # Recommendation Section
    if recommendation_lines:
        elements.append(Spacer(1, 12))
        elements.append(Paragraph("<b>✅ "+t('best_reco')+"</b>", styles["Heading3"]))
        for line in recommendation_lines:
            elements.append(Paragraph(line, highlight_style))

    doc.build(elements)
    buffer.seek(0)

    with open("static/hardware_quotation.pdf", "wb") as f:
        f.write(buffer.getvalue())
    return buffer


# ────────────────────────────────────────────────────────────────────────────────
# Faster Whisper STT + gTTS TTS Helpers
# ────────────────────────────────────────────────────────────────────────────────
@st.cache_resource(show_spinner="Loading Faster Whisper model…")
def load_whisper_local(model_name: str = "base"):
    device = "cuda" if torch.cuda.is_available() else "cpu"
    compute_type = "int8_float16" if device=="cuda" else "int8"
    return WhisperModel(model_name, device=device, compute_type=compute_type, local_files_only=False)

def transcribe_audio(wav_input) -> str:
    wav_bytes = wav_input.get("bytes", b"") if isinstance(wav_input, dict) else wav_input or b""
    if not wav_bytes:
        return ""
    with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as tmp:
        tmp.write(wav_bytes); path = tmp.name
    try:
        lang_code = st.session_state.get("lang_code", "en")
        segs, _ = load_whisper_local().transcribe(path, beam_size=5,
                                                  suppress_blank=True,
                                                  language=lang_code)
        return "".join(s.text for s in segs).strip()
    finally:
        os.remove(path)

def _preprocess_for_tts(raw: str) -> str:
    lines = []
    for line in raw.splitlines():
        line = line.strip().lstrip("#").lstrip("-*•").strip()
        if line:
            lines.append(line)
    return ".  ".join(lines)

def tts_stream(text: str) -> BytesIO:
    code = st.session_state.get("lang_code", "en")
    tld  = "co.jp" if code == "ja" else "co.uk"
    mp3  = BytesIO()
    gTTS(text=_preprocess_for_tts(text), lang=code, tld=tld, slow=False).write_to_fp(mp3)
    mp3.seek(0)
    return mp3

# def speak_async(text:str)->None:
#     threading.Thread(target=lambda: speaker.Speak(text),daemon=True).start()


# --- LangGraph App Logic ---
class State(TypedDict):
    question: str
    context: List[Document]
    answer: str
    feedback: str

def retrieve(state: State):
    print(state["question"])
    if state.get("feedback"):
        state["question"] += f"\nAdditional feedback: {state['feedback']}"
    print(vector_store.similarity_search(state["question"]))
    return {"context": vector_store.similarity_search(state["question"])}

def generate(state: State):
    docs_content = "\n\n".join(doc.page_content for doc in state["context"])

    question_with_feedback = state["question"]
    if state.get("feedback"):
        question_with_feedback += f"\nAdditional Feedback: {state['feedback']}"
        print(question_with_feedback)
        state["question"]= question_with_feedback

    prompt = build_prompt()
    messages = prompt.invoke({"question": state["question"], "context": docs_content})
    print(state["question"])
    response = llm.invoke(messages)
    return {"answer": response.content}

def extract_recommendation_text(response: str) -> str:
    # This will match 'Recommendation' followed by optional colon and any whitespace, then capture the rest
    match = re.search(r"Recommendation[:\s]*([\s\S]+)", response, re.IGNORECASE)
    if match:
        return match.group(1).strip()
    return ""


graph_builder = StateGraph(State).add_sequence([retrieve, generate])
graph_builder.add_edge(START, "retrieve")
graph = graph_builder.compile()

# ---------- Constants ----------
LOG_DIR = "chat_logs"
os.makedirs(LOG_DIR, exist_ok=True)

# ---------- Load logo ----------
def get_base64_logo(image_path: str) -> str:
    with open(image_path, "rb") as f:
        return base64.b64encode(f.read()).decode("utf-8")

logo_base64 = get_base64_logo("otsuka_im.png")  # Local image path

# ---------- CSS ----------
st.markdown(
    f"""
    <style>
        body {{
            background-color: #ffedec;
        }}
        .main-title h1 {{
            color: #002060;
        }}
        .client-box {{
            background-color: #ffffff;
            padding: 1.5rem;
            border-radius: 12px;
            box-shadow: 0 4px 10px rgba(0,0,0,0.1);
            margin-bottom: 2rem;
        }}
        .stTextInput > label {{
            font-weight: bold;
            color: #002060;
        }}
        div[data-testid="stHeader"] {{
            background-color: white;
        }}
        .logo-container {{
            position: absolute;
            top: 15px;
            right: 15px;
            z-index: 1000;
        }}
    </style>
    <div class="logo-container">
        <img src="data:image/png;base64,{logo_base64}" width="130">
    </div>
    """,
    unsafe_allow_html=True
)

# ---------- Title ----------
st.markdown('<div class="main-title">', unsafe_allow_html=True)
# st.title("💻 Computer Hardware Sales Assistant")
st.title(t("title"))
st.markdown('</div>', unsafe_allow_html=True)

# ---------- Language switcher ----------
selected_lang = st.sidebar.radio("Language / 言語", list(LANGUAGES.keys()))
if "lang_code" not in st.session_state or st.session_state.lang_code != LANGUAGES[selected_lang]:
    st.session_state.lang_code = LANGUAGES[selected_lang]
    st.rerun()


# ---------- Sidebar: Sessions ----------
st.sidebar.header(t("sidebar_sessions"))
session_files = [f for f in os.listdir(LOG_DIR) if f.endswith(".csv")]
selected_session = st.sidebar.selectbox(t("load_session"), [""] + session_files)

if selected_session:
    df = pd.read_csv(os.path.join(LOG_DIR, selected_session))
    st.sidebar.markdown(t("loaded_data_hdr"))
    st.sidebar.dataframe(df)

# ---------- Initialize session state ----------
# for key in ["result", "feedback", "user_query", "active", "pdf_bytes", "slide_bytes", "chat_history"]:
#     if key not in st.session_state:
#         st.session_state[key] = None if key != "chat_history" else []

# ---------- Initialize session state ----------
for key in [
    "result", "feedback", "user_query", "active",
    "pdf_bytes", "slide_bytes", "chat_history",
    "tts_mp3"          # <─---- NEW
]:
    if key not in st.session_state:
        st.session_state[key] = None if key not in ["chat_history","tts_mp3"] else ([] if key=="chat_history" else None)


# helper just before the expander, keeps the emojis
def ci(lbl_en: str, lbl_ja: str) -> str:
    return _l(lbl_en, lbl_ja)

with st.expander(t("client_expander"), expanded=True):
    col1, col2 = st.columns(2)
    with col1:
        client_name  = st.text_input( f"👤 {ci('Client Name','お名前')}", placeholder=t("name_ph") )
        client_email = st.text_input( f"📧 {ci('Email','メール')}",                         placeholder=t("email_ph") )
    with col2:
        client_company = st.text_input( f"🏢 {ci('Company','会社名')}",    placeholder=t("company_ph") )
        client_phone   = st.text_input( f"📞 {ci('Phone','電話')}",        placeholder=t("phone_ph") )

    st.session_state.client_info = {"name": client_name, "company": client_company,
                                    "email": client_email, "phone": client_phone}



# Initialize session state for STT/TTS
for key in ["text_query", "just_transcribed", "tts_play", "tts_mp3","feedback_text", "just_transcribed_fb"]:
    if key not in st.session_state:
        st.session_state[key] = False if key in ["just_transcribed","tts_play","just_transcribed_fb"] else None


# ---------- Query Input ----------
st.divider()
# st.header("How can I help you today? 😊")
st.header(t("help_header"))

# microphone
audio_main = mic_recorder(start_prompt="🎤", stop_prompt="⏹", key="recorder_main")
if audio_main and not st.session_state.just_transcribed:
    with st.spinner(t("transcribing")):
        spoken = transcribe_audio(audio_main)
        if spoken:
            st.session_state.text_query = spoken
            st.session_state.just_transcribed = True


if "trigger_sample_query" not in st.session_state:
    st.session_state.trigger_sample_query = False
if "selected_sample_query" not in st.session_state:
    st.session_state.selected_sample_query = ""

# Sample query buttons
col1, col2, col3 = st.columns(3)
with col1:
    if st.button(t("sample1")):
        st.session_state.selected_sample_query = "Suggest me a lightweight business laptop"
        st.session_state.trigger_sample_query = True
        st.rerun()
with col2:
    if st.button(t("sample2")):
        st.session_state.selected_sample_query = "Suggest me a good gaming mouse"
        st.session_state.trigger_sample_query = True
        st.rerun()
with col3:
    if st.button(t("sample3")):
        st.session_state.selected_sample_query = "Help me with a 4K webcam for making reel shots"
        st.session_state.trigger_sample_query = True
        st.rerun()



# user_query = st.text_input("Enter your query:", placeholder="E.g., Best PC setup for video editing...")

# user_query = st.text_input(
#     "Enter your query:",
#     value=st.session_state.selected_sample_query,
#     placeholder="E.g., Best PC setup for video editing…",
#     key = "text_query"
# )

user_query = st.text_input(
    t("query_placeholder"),
    value=st.session_state.selected_sample_query,
    placeholder=t("query_placeholder"),
    key="text_query"
)

# Handle query trigger from sample button
if st.session_state.trigger_sample_query and st.session_state.selected_sample_query:
    with st.spinner(t("finding")):
        st.session_state.user_query = st.session_state.selected_sample_query
        st.session_state.result = graph.invoke({"question": st.session_state.user_query})
        recom = extract_recommendation_text(st.session_state.result["answer"])
        tts_bytes_obj = tts_stream(recom)
        st.session_state.tts_mp3 = tts_bytes_obj
        st.session_state.active = True
        st.session_state.feedback = None
        st.session_state.chat_history.append({
            "message": st.session_state.user_query,
            "answer": st.session_state.result["answer"],
            "feedback": ""
        })
        st.session_state.trigger_sample_query = False  # Reset the trigger

# ---------- On Query Submit ----------
if user_query and st.button(t("get_reco")):
    with st.spinner(t("finding")):
        st.session_state.user_query = user_query
        st.session_state.result = graph.invoke({"question": user_query})  # <-- Replace with your model
        recom = extract_recommendation_text(st.session_state.result["answer"])
        tts_bytes_obj = tts_stream(recom)
        st.session_state.tts_mp3 = tts_bytes_obj
        st.session_state.active = True
        st.session_state.feedback = None

        # Save to chat history
        st.session_state.chat_history.append({
            "message": user_query,
            "answer": st.session_state.result["answer"],
            "feedback": ""
        })

# ---------- Result Display ----------
if st.session_state.result:
    st.subheader(t("quotation_header"))
    ans = st.session_state.result["answer"]
    st.write(st.session_state.result["answer"])

    # Token stats
    enc = tiktoken.encoding_for_model("gpt-4")
    input_tokens = len(enc.encode(st.session_state.user_query))
    output_tokens = len(enc.encode(st.session_state.result["answer"]))
    st.markdown(f"🔢 Input tokens: {input_tokens} | Output tokens: {output_tokens} | Total: {input_tokens + output_tokens}")

    # Export PDFs and Slides
    st.session_state.pdf_bytes = generate_pdf(st.session_state.result["answer"], st.session_state.client_info)
    st.session_state.slide_bytes = generate_slides(st.session_state.result["answer"], st.session_state.client_info)

    # Export options
    st.markdown(t("export_header"))
    col1, col2, col3 = st.columns(3)
    with col1:
        st.download_button(t("export_pdf"), data=st.session_state.pdf_bytes, file_name="hardware_quotation.pdf", mime="application/pdf")
    with col2:
        st.download_button(t("export_ppt"), data=st.session_state.slide_bytes, file_name="hardware_quotation.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    with col3:
        if st.button(t("email_quote")):
            if not client_email:
                st.error(t("email_err"))
            else:
                with st.spinner(t("sending_email")):
                    success, msg = send_email_with_attachment(
                        to_email=client_email,
                        subject="Hardware Quotation from Otsuka Shokai",
                        body=f"Dear {client_name},\n\nPlease find attached your hardware quotation.\n\nRegards,\nOtsuka Shokai"
                    )
                    st.success(msg) if success else st.error(msg)


    # if st.button("🔊 Speak this quotation"):
    #     with st.spinner("Generating speech…"):
    #         recom=extract_recommendation_text(ans)
    #         print(f"Ans {recom}")
    #         mp3 = tts_stream(recom)
    #         st.audio(mp3.read(), format="audio/mp3")

    # ----- Speak button -----
    if st.button(t("speak_btn")):
        if st.session_state.tts_mp3 is None:
            st.warning(t("audio_wait"))
        else:
            st.audio(st.session_state.tts_mp3, format="audio/mp3")



    # ---------- Feedback Loop ----------
    if st.session_state.active:
        st.markdown("---")
        st.subheader(t("feedback_header"))



        audio_fb = mic_recorder(start_prompt="🎤", stop_prompt="⏹", key="recorder_fb")
        if audio_fb and not st.session_state.just_transcribed_fb:
            with st.spinner(t("transcribing_fb")):
                fb_spoken = transcribe_audio(audio_fb)
                if fb_spoken:
                    st.session_state.feedback_text = fb_spoken
                    st.session_state.just_transcribed_fb = True

        feedback_text = st.text_area(
            t("feedback_placeholder"),
            value=st.session_state.feedback_text,
            key="feedback_text",
            # on_change=lambda: st.session_state.setitem("just_transcribed_fb", False)
        )



        # feedback_text = st.text_area("Suggest changes or type 'thanks' to finalize:")

        if st.button(t("submit_feedback")):
            if feedback_text:
                if "thank" in feedback_text.lower():
                    st.success(t("thank_finalized"))

                    st.session_state.active = False
                    # ---------- Save Session ----------
                    

# Add this near the bottom of your script (e.g., after "Save Session" button)
                else:
                    with st.spinner(t("incorporating")):
                        st.session_state.feedback = feedback_text
                        st.session_state.user_query += f"\n {feedback_text}"
                        #st.session_state.user_query+= f"\n {feedback_text}"
                        st.session_state.result = graph.invoke({
                            "question": st.session_state.user_query,
                            "feedback": feedback_text
                        })

                        recom = extract_recommendation_text(st.session_state.result["answer"])
                        tts_bytes_obj = tts_stream(recom)
                        st.session_state.tts_mp3 = tts_bytes_obj

                        # Save revised answer
                        st.session_state.chat_history.append({
                            "message": st.session_state.user_query,
                            "answer": st.session_state.result["answer"],
                            "feedback": feedback_text
                        })

                    st.subheader(t("revised_quotation"))
                    st.write(st.session_state.result["answer"])
                    input_tokens = len(enc.encode(st.session_state.user_query))
                    output_tokens = len(enc.encode(st.session_state.result["answer"]))
                    st.markdown(f"🔢 Input tokens: {input_tokens} | Output tokens: {output_tokens} | Total: {input_tokens + output_tokens}")
                    st.session_state.pdf_bytes = generate_pdf(st.session_state.result["answer"], st.session_state.client_info)
                    st.session_state.slide_bytes = generate_slides(st.session_state.result["answer"], st.session_state.client_info)
                    st.rerun()
if st.button(t("save_chat")):
    if st.session_state.chat_history:
        timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
        name_user= st.session_state.client_info["name"]
        filename = f"{LOG_DIR}/session_{name_user}.csv"
        pd.DataFrame(st.session_state.chat_history).to_csv(filename, index=False)
        st.success(t("session_saved").format(file=filename))
        st.rerun()
    else:
        st.warning(t("no_chat_warning"))

if st.button(t("restart")):
    # Clear all session state variables relevant to this chat
    keys_to_clear = [
        "result", "feedback", "user_query", "active", 
        "pdf_bytes", "slide_bytes", "chat_history", 
        "client_info", "tts_mp3"
    ]
    for key in keys_to_clear:
        if key in st.session_state:
            del st.session_state[key]
    st.rerun()